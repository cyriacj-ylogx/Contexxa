"""Document loaders extracted from rag_responder for the PageIndex branch.

Standalone copy of HelpCenterAgent's file-loading subtree so the vectorless
pipeline can load .docx/.pdf/.txt/.md/.csv/.json/.html/.xlsx/.pptx files
WITHOUT importing chromadb / sentence-transformers / torch. Behaviour is
identical to rag_responder's loaders (same block metadata: source,
document_type, section_path, block_type, heading flags).
"""

import hashlib
import json
import os
import re
from typing import Dict, List, Sequence, Tuple

from docx import Document as DocxDocument
from docx.document import Document as DocxNativeDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from dataclasses import dataclass, field
import pypdf


@dataclass
class Document:
    """Minimal Document container (replaces langchain.schema.Document)."""
    page_content: str
    metadata: dict = field(default_factory=dict)


class DocLoaders:
    """File -> List[Document] block loaders. All methods are class/static."""

    MIN_BODY_WORDS = 7
    MIN_BODY_CHARS = 45
    TOC_PAGE_SCAN_LIMIT = 4
    MAX_TEXT_FILE_BYTES = 10 * 1024 * 1024  # 10 MB guard for plain-text loads

    GENERIC_TEXT_EXTENSIONS = (
        ".txt",
        ".md",
        ".markdown",
        ".rst",
        ".csv",
        ".tsv",
        ".json",
        ".html",
        ".htm",
        ".log",
        ".text",
        ".rtf",
    )

    # ── small text helpers ────────────────────────────────────────────
    @staticmethod
    def _normalize_text(text: str) -> str:
        return re.sub(r"\s+", " ", text or "").strip()

    @staticmethod
    def _file_key(source_path: str) -> str:
        """Stable short key for a source file (same scheme as rag_responder)."""
        basename = os.path.basename(source_path or "")
        return hashlib.sha1(basename.encode("utf-8")).hexdigest()[:8]

    @staticmethod
    def _is_heading(text: str, style_name: str = "") -> bool:
        normalized = DocLoaders._normalize_text(text)
        if not normalized:
            return False
        if style_name.lower().startswith("heading"):
            return True
        if len(normalized) <= 120:
            if re.match(r"^\d+(\.\d+)*\s+[A-Za-z].*", normalized):
                return True
            if re.match(r"^[A-Z][A-Za-z/&,\-\s\(\)]+?\s+\d+$", normalized):
                return True
            if normalized.isupper() and len(normalized.split()) <= 12:
                return True
            if (
                normalized[0].isupper()
                and not normalized.endswith(".")
                and len(normalized.split()) <= 10
            ):
                return True
        return False

    @staticmethod
    def _heading_level(style_name: str) -> int:
        match = re.search(r"heading\s*(\d+)", style_name.lower())
        if match:
            return max(1, min(6, int(match.group(1))))
        if style_name.lower().startswith("heading"):
            return 2
        return 1

    @staticmethod
    def _is_bullet_line(text: str) -> bool:
        return bool(
            re.match(r"^([-\*]|\d+[\.\)])\s+", DocLoaders._normalize_text(text))
        )

    @staticmethod
    def _is_table_like(text: str) -> bool:
        normalized = DocLoaders._normalize_text(text)
        if "|" in normalized:
            return True
        return bool(re.search(r"\s{2,}", normalized) and len(normalized.split()) <= 18)

    @staticmethod
    def _looks_like_toc_line(text: str) -> bool:
        normalized = DocLoaders._normalize_text(text)
        return bool(
            normalized
            and len(normalized.split()) <= 10
            and re.match(r"^[A-Za-z][A-Za-z/&,\-\s\(\)]+?\s+\d+$", normalized)
        )

    @classmethod
    def _is_probable_pdf_toc_page(cls, lines: Sequence[str], page_number: int) -> bool:
        try:
            limit = int(os.environ.get("TOC_PAGE_SCAN_LIMIT", cls.TOC_PAGE_SCAN_LIMIT))
        except ValueError:
            limit = int(cls.TOC_PAGE_SCAN_LIMIT)
        if page_number + 1 > limit:
            return False
        normalized = [cls._normalize_text(line) for line in lines if cls._normalize_text(line)]
        if len(normalized) < 8:
            return False
        toc_lines = sum(1 for line in normalized if cls._looks_like_toc_line(line))
        long_lines = sum(1 for line in normalized if len(line.split()) >= 12)
        return toc_lines >= 8 and toc_lines >= max(1, len(normalized) // 2) and long_lines <= 3

    @classmethod
    def _classify_block(cls, text: str, style_name: str = "") -> str:
        if cls._is_heading(text, style_name):
            return "heading"
        if cls._is_bullet_line(text):
            return "bullet"
        if cls._is_table_like(text):
            return "table_like"
        return "paragraph"

    @staticmethod
    def _iter_docx_items(doc: DocxNativeDocument):
        for child in doc.element.body.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, doc)
            elif isinstance(child, CT_Tbl):
                yield Table(child, doc)

    @classmethod
    def _docx_section_path(cls, section_stack: Sequence[Tuple[int, str]]) -> str:
        return " > ".join(title for _, title in section_stack)

    @classmethod
    def _build_metadata(
        cls,
        source: str,
        document_type: str,
        section_stack: Sequence[Tuple[int, str]],
        extra: Dict[str, object] = None,
    ) -> Dict[str, object]:
        metadata: Dict[str, object] = {
            "source": source,
            "document_type": document_type,
            "section_path": cls._docx_section_path(section_stack),
            "section_depth": len(section_stack),
        }
        if extra:
            metadata.update(extra)
        return metadata

    # ── HTML / generic text ───────────────────────────────────────────
    @staticmethod
    def _strip_html(raw: str) -> str:
        without_blocks = re.sub(
            r"<(script|style)[^>]*>.*?</\1>", " ", raw, flags=re.IGNORECASE | re.DOTALL
        )
        with_breaks = re.sub(
            r"</(p|div|li|tr|h[1-6]|br|section|article)\s*>",
            "\n",
            without_blocks,
            flags=re.IGNORECASE,
        )
        with_breaks = re.sub(r"<br\s*/?>", "\n", with_breaks, flags=re.IGNORECASE)
        no_tags = re.sub(r"<[^>]+>", " ", with_breaks)
        replacements = {
            "&nbsp;": " ",
            "&amp;": "&",
            "&lt;": "<",
            "&gt;": ">",
            "&quot;": '"',
            "&#39;": "'",
            "&apos;": "'",
        }
        for entity, char in replacements.items():
            no_tags = no_tags.replace(entity, char)
        return no_tags

    @classmethod
    def _text_to_blocks(
        cls, file_path: str, text: str, source_type: str, merge_short_lines: bool = False
    ) -> List[Document]:
        blocks: List[Document] = []
        section_stack: List[Tuple[int, str]] = []
        short_line_buffer: List[str] = []
        doc_title: List[str] = []

        def flush_short_buffer() -> None:
            if not short_line_buffer:
                return
            title_part = doc_title[0] if doc_title else ""
            section_part = (
                cls._normalize_text(section_stack[-1][1]) if section_stack else ""
            )
            if title_part and title_part != section_part:
                prefix = f"{title_part} {section_part}: "
            elif section_part:
                prefix = f"{section_part}: "
            else:
                prefix = ""
            merged = cls._normalize_text(prefix + " ".join(short_line_buffer))
            short_line_buffer.clear()
            if not merged:
                return
            blocks.append(
                Document(
                    page_content=merged,
                    metadata=cls._build_metadata(
                        file_path,
                        source_type,
                        section_stack,
                        {"block_type": "bullet"},
                    ),
                )
            )

        raw_paragraphs = re.split(r"\n\s*\n", text)
        for raw_paragraph in raw_paragraphs:
            for raw_line in raw_paragraph.split("\n"):
                stripped = raw_line.strip()
                if not stripped:
                    continue

                md_heading = re.match(r"^(#{1,6})\s+(.*)$", stripped)
                if md_heading:
                    flush_short_buffer()
                    level = len(md_heading.group(1))
                    title = cls._normalize_text(md_heading.group(2))
                    if title:
                        if not doc_title:
                            doc_title.append(title)
                        section_stack = [e for e in section_stack if e[0] < level]
                        section_stack.append((level, title))
                        blocks.append(
                            Document(
                                page_content=title,
                                metadata=cls._build_metadata(
                                    file_path,
                                    source_type,
                                    section_stack,
                                    {
                                        "block_type": "heading",
                                        "heading_level": level,
                                        "is_heading_only": True,
                                    },
                                ),
                            )
                        )
                    continue

                normalized = cls._normalize_text(stripped)
                if not normalized:
                    continue

                if merge_short_lines and len(normalized) < cls.MIN_BODY_CHARS:
                    short_line_buffer.append(normalized)
                    continue

                if merge_short_lines and normalized.endswith(":"):
                    flush_short_buffer()
                    short_line_buffer.append(normalized)
                    continue

                if cls._is_heading(normalized):
                    flush_short_buffer()
                    level = 2
                    if not doc_title:
                        doc_title.append(normalized)
                    section_stack = [e for e in section_stack if e[0] < level]
                    section_stack.append((level, normalized))
                    blocks.append(
                        Document(
                            page_content=normalized,
                            metadata=cls._build_metadata(
                                file_path,
                                source_type,
                                section_stack,
                                {
                                    "block_type": "heading",
                                    "heading_level": level,
                                    "is_heading_only": True,
                                },
                            ),
                        )
                    )
                    continue

                flush_short_buffer()
                blocks.append(
                    Document(
                        page_content=normalized,
                        metadata=cls._build_metadata(
                            file_path,
                            source_type,
                            section_stack,
                            {"block_type": cls._classify_block(normalized)},
                        ),
                    )
                )

        flush_short_buffer()
        return blocks

    @classmethod
    def _flatten_json_to_blocks(
        cls, file_path: str, obj: object, _section_stack: List[Tuple[int, str]] = None
    ) -> List[Document]:
        """Recursively flatten parsed JSON into readable Document blocks."""
        if _section_stack is None:
            _section_stack = []
        blocks: List[Document] = []

        def make_doc(text: str, stack: List[Tuple[int, str]]) -> None:
            normalized = cls._normalize_text(text)
            if not normalized:
                return
            blocks.append(
                Document(
                    page_content=normalized,
                    metadata=cls._build_metadata(
                        file_path, "json", stack, {"block_type": "paragraph"}
                    ),
                )
            )

        def flatten_value(v: object) -> str:
            if isinstance(v, str):
                return cls._normalize_text(v)
            if isinstance(v, (int, float, bool)):
                return str(v)
            if isinstance(v, list):
                parts = [flatten_value(item) for item in v if not isinstance(item, (dict, list))]
                return ", ".join(p for p in parts if p)
            return ""

        def process(node: object, stack: List[Tuple[int, str]]) -> None:
            if isinstance(node, dict):
                name_key = next(
                    (k for k in ("name", "title", "id", "type") if k in node), None
                )
                if name_key:
                    heading = cls._normalize_text(str(node[name_key]))
                    new_stack = [e for e in stack if e[0] < 2] + [(2, heading)]
                else:
                    new_stack = stack

                parts: List[str] = []
                for k, v in node.items():
                    if isinstance(v, (dict, list)):
                        continue
                    val = flatten_value(v)
                    if val:
                        parts.append(f"{cls._normalize_text(str(k))}: {val}")
                if parts:
                    make_doc(" | ".join(parts), new_stack)

                for v in node.values():
                    if isinstance(v, (dict, list)):
                        process(v, new_stack)

            elif isinstance(node, list):
                for item in node:
                    process(item, stack)

            elif isinstance(node, str):
                val = cls._normalize_text(node)
                if val and len(val) >= cls.MIN_BODY_CHARS:
                    make_doc(val, stack)

        process(obj, _section_stack)
        return blocks

    # ── per-format loaders ────────────────────────────────────────────
    @classmethod
    def _load_text(cls, file_path):
        """Generic loader for text-extractable files (txt/md/csv/json/html/...)."""
        try:
            try:
                if os.path.getsize(file_path) > cls.MAX_TEXT_FILE_BYTES:
                    print(f"Skipping oversized text file: {file_path}")
                    return []
            except OSError:
                pass

            with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
                raw = handle.read()

            ext = os.path.splitext(file_path)[1].lower()

            if ext in (".html", ".htm"):
                text = cls._strip_html(raw)
            elif ext in (".csv", ".tsv"):
                delimiter = "\t" if ext == ".tsv" else ","
                rows = []
                for line in raw.splitlines():
                    if not line.strip():
                        continue
                    cells = [cell.strip() for cell in line.split(delimiter)]
                    row_text = " | ".join(cell for cell in cells if cell)
                    if row_text:
                        rows.append(row_text)
                header = rows[0] if rows else ""
                data_rows = rows[1:] if len(rows) > 1 else rows
                CSV_CHUNK = 5
                csv_section: List[Tuple[int, str]] = [(1, header)] if header else []
                csv_blocks: List[Document] = []
                for _i in range(0, max(len(data_rows), 1), CSV_CHUNK):
                    group = data_rows[_i: _i + CSV_CHUNK]
                    chunk_text = cls._normalize_text(
                        (header + "\n" if header else "") + "\n".join(group)
                    )
                    if not chunk_text:
                        continue
                    csv_blocks.append(
                        Document(
                            page_content=chunk_text,
                            metadata=cls._build_metadata(
                                file_path,
                                "csv",
                                csv_section,
                                {"block_type": "table_row"},
                            ),
                        )
                    )
                return csv_blocks
            elif ext == ".json":
                try:
                    parsed = json.loads(raw)
                    json_blocks = cls._flatten_json_to_blocks(file_path, parsed)
                    if json_blocks:
                        return json_blocks
                    text = json.dumps(parsed, indent=2, ensure_ascii=False)
                except Exception:
                    text = raw
            else:
                text = raw

            source_type = {
                ".csv": "csv",
                ".tsv": "csv",
                ".json": "json",
                ".html": "html",
                ".htm": "html",
                ".md": "markdown",
                ".markdown": "markdown",
            }.get(ext, "text")

            merge_bullets = ext in (".html", ".htm")
            return cls._text_to_blocks(
                file_path, text, source_type, merge_short_lines=merge_bullets
            )
        except Exception as e:
            raise Exception(f"Failed to load text file: {e}")

    @classmethod
    def _load_html(cls, file_path):
        """Parse HTML using BeautifulSoup to preserve h1/h2/h3 section structure."""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                raw = f.read()
            soup = BeautifulSoup(raw, "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()

            source = os.path.basename(file_path)
            blocks: List[Document] = []
            section_stack: List[Tuple[int, str]] = []

            def flush_element(el, stack):
                text = cls._normalize_text(el.get_text(" ", strip=True))
                if len(text) < cls.MIN_BODY_WORDS:
                    return
                blocks.append(Document(
                    page_content=text,
                    metadata=cls._build_metadata(file_path, "html", stack, {"block_type": "paragraph"}),
                ))

            HEADING_TAGS = {"h1": 1, "h2": 2, "h3": 3, "h4": 4}
            for el in soup.find_all(True):
                if el.name in HEADING_TAGS:
                    level = HEADING_TAGS[el.name]
                    title = cls._normalize_text(el.get_text(" ", strip=True))
                    if not title:
                        continue
                    section_stack = [e for e in section_stack if e[0] < level]
                    section_stack.append((level, title))
                    blocks.append(Document(
                        page_content=title,
                        metadata=cls._build_metadata(file_path, "html", section_stack,
                                                     {"block_type": "heading", "heading_level": level}),
                    ))
                elif el.name in ("p", "li") and el.find_parent(HEADING_TAGS.keys()) is None:
                    # avoid double-processing nested elements inside headings
                    if el.find_parent(["p", "li"]) is None or el.name == "li":
                        flush_element(el, list(section_stack))

            return blocks
        except ImportError:
            pass
        except Exception as e:
            raise Exception(f"Failed to load html file: {e}")
        # fallback to text loader without merge_short_lines
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read()
        text = cls._strip_html(raw)
        return cls._text_to_blocks(file_path, text, "html", merge_short_lines=False)

    @classmethod
    def _load_office(cls, file_path):
        """Optional loader for .pptx/.xlsx — only if the libs are installed."""
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".pptx":
                from pptx import Presentation  # type: ignore

                prs = Presentation(file_path)
                lines: List[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False):
                            for paragraph in shape.text_frame.paragraphs:
                                line = cls._normalize_text(
                                    "".join(run.text for run in paragraph.runs)
                                )
                                if line:
                                    lines.append(line)
                return cls._text_to_blocks(file_path, "\n".join(lines), "pptx")
            if ext == ".xlsx":
                from openpyxl import load_workbook  # type: ignore

                wb = load_workbook(file_path, read_only=True, data_only=True)
                rows: List[str] = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        cells = [str(c).strip() for c in row if c is not None]
                        row_text = " | ".join(cell for cell in cells if cell)
                        if row_text:
                            rows.append(row_text)
                return cls._text_to_blocks(file_path, "\n".join(rows), "xlsx")
        except ImportError:
            print(f"Skipping {file_path}: optional library for {ext} not installed")
            return []
        except Exception as e:
            print(f"Unexpected error loading {file_path}: {e}")
            return []
        return []

    @classmethod
    def _load_docx(cls, file_path):
        try:
            doc = DocxDocument(file_path)
            blocks = []
            section_stack: List[Tuple[int, str]] = []
            table_index = 0

            for item in cls._iter_docx_items(doc):
                if isinstance(item, Paragraph):
                    text = cls._normalize_text(item.text)
                    if not text:
                        continue
                    style_name = getattr(getattr(item, "style", None), "name", "") or ""

                    if cls._is_heading(text, style_name):
                        level = cls._heading_level(style_name)
                        section_stack = [
                            entry for entry in section_stack if entry[0] < level
                        ]
                        section_stack.append((level, text))
                        blocks.append(
                            Document(
                                page_content=text,
                                metadata=cls._build_metadata(
                                    file_path,
                                    "docx",
                                    section_stack,
                                    {
                                        "block_type": "heading",
                                        "heading_level": level,
                                        "is_heading_only": True,
                                    },
                                ),
                            )
                        )
                        continue

                    blocks.append(
                        Document(
                            page_content=text,
                            metadata=cls._build_metadata(
                                file_path,
                                "docx",
                                section_stack,
                                {
                                    "block_type": cls._classify_block(text, style_name),
                                    "paragraph_style": style_name,
                                },
                            ),
                        )
                    )

                elif isinstance(item, Table):
                    table_index += 1
                    for row_index, row in enumerate(item.rows, start=1):
                        row_content = [cls._normalize_text(cell.text) for cell in row.cells]
                        row_text = cls._normalize_text(
                            " | ".join(part for part in row_content if part)
                        )
                        if not row_text:
                            continue
                        blocks.append(
                            Document(
                                page_content=row_text,
                                metadata=cls._build_metadata(
                                    file_path,
                                    "docx",
                                    section_stack,
                                    {
                                        "block_type": "table_row",
                                        "table_index": table_index,
                                        "row_index": row_index,
                                    },
                                ),
                            )
                        )

            return blocks
        except KeyError as e:
            raise KeyError(
                f"There is no item named '[Content_Types].xml' in the archive: {e}"
            )
        except Exception as e:
            raise Exception(f"Failed to load .docx file: {e}")

    @classmethod
    def _split_pdf_page_into_blocks(
        cls,
        file_path: str,
        page_doc: Document,
        page_number: int,
        section_stack: Sequence[Tuple[int, str]],
    ):
        raw_lines = page_doc.page_content.splitlines()
        if cls._is_probable_pdf_toc_page(raw_lines, page_number):
            return [], list(section_stack)

        blocks = []
        current_sections: List[Tuple[int, str]] = list(section_stack)
        paragraph_buffer: List[str] = []
        paragraph_index = 0

        def flush_paragraph_buffer():
            nonlocal paragraph_index
            if not paragraph_buffer:
                return
            paragraph_index += 1
            paragraph_text = cls._normalize_text(" ".join(paragraph_buffer))
            paragraph_buffer.clear()
            if not paragraph_text:
                return
            blocks.append(
                Document(
                    page_content=paragraph_text,
                    metadata=cls._build_metadata(
                        file_path,
                        "pdf",
                        current_sections,
                        {
                            "block_type": "paragraph",
                            "page_number": page_number + 1,
                            "paragraph_index": paragraph_index,
                        },
                    ),
                )
            )

        for raw_line in raw_lines:
            line = cls._normalize_text(raw_line)
            if not line:
                flush_paragraph_buffer()
                continue

            if cls._is_heading(line):
                flush_paragraph_buffer()
                level = cls._heading_level("")
                current_sections = [entry for entry in current_sections if entry[0] < level]
                current_sections.append((level, line))
                blocks.append(
                    Document(
                        page_content=line,
                        metadata=cls._build_metadata(
                            file_path,
                            "pdf",
                            current_sections,
                            {
                                "block_type": "heading",
                                "page_number": page_number + 1,
                                "heading_level": level,
                                "is_heading_only": True,
                            },
                        ),
                    )
                )
                continue

            if cls._is_bullet_line(line):
                flush_paragraph_buffer()
                blocks.append(
                    Document(
                        page_content=line,
                        metadata=cls._build_metadata(
                            file_path,
                            "pdf",
                            current_sections,
                            {"block_type": "bullet", "page_number": page_number + 1},
                        ),
                    )
                )
                continue

            if cls._is_table_like(line):
                flush_paragraph_buffer()
                blocks.append(
                    Document(
                        page_content=line,
                        metadata=cls._build_metadata(
                            file_path,
                            "pdf",
                            current_sections,
                            {"block_type": "table_like", "page_number": page_number + 1},
                        ),
                    )
                )
                continue

            paragraph_buffer.append(line)

        flush_paragraph_buffer()
        return blocks, current_sections

    @classmethod
    def _load_pdf(cls, file_path):
        try:
            reader = pypdf.PdfReader(file_path)
            blocks = []
            section_stack: List[Tuple[int, str]] = []
            for page_number, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                page_doc = Document(
                    page_content=text,
                    metadata={"source": file_path, "page": page_number},
                )
                page_blocks, section_stack = cls._split_pdf_page_into_blocks(
                    file_path, page_doc, page_number, section_stack
                )
                blocks.extend(page_blocks)
            return blocks
        except Exception as e:
            raise Exception(f"Failed to load .pdf file: {e}")

    # ── public entry points ───────────────────────────────────────────
    @classmethod
    def load_single_file(cls, file_path: str) -> List[Document]:
        """Route one file to the right loader by extension."""
        lower = file_path.lower()
        ext = os.path.splitext(lower)[1]
        if lower.endswith(".docx"):
            return cls._load_docx(file_path)
        if lower.endswith(".pdf"):
            return cls._load_pdf(file_path)
        if ext in (".html", ".htm"):
            return cls._load_html(file_path)
        if ext in cls.GENERIC_TEXT_EXTENSIONS:
            return cls._load_text(file_path)
        if ext in (".pptx", ".xlsx"):
            return cls._load_office(file_path)
        print(f"Skipping unsupported file type: {file_path}")
        return []

    @classmethod
    def load_docs(cls, directory: str) -> List[Document]:
        documents = []
        for filename in sorted(os.listdir(directory)):
            doc_path = os.path.join(directory, filename)
            if not os.path.isfile(doc_path):
                continue
            try:
                documents.extend(cls.load_single_file(doc_path))
            except Exception as e:
                print(f"Unexpected error loading {doc_path}: {e}")
        return documents

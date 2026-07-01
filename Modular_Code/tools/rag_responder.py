import hashlib
import json
import math
import os
import re
from collections import Counter
from typing import Dict, List, Sequence, Tuple

import numpy as np
from docx import Document as DocxDocument
from docx.document import Document as DocxNativeDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from langchain.chains.question_answering import load_qa_chain
from langchain.chat_models import AzureChatOpenAI
from langchain.document_loaders import PyPDFLoader
from langchain.embeddings import SentenceTransformerEmbeddings
from langchain.prompts import PromptTemplate
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma
from sentence_transformers import CrossEncoder, SentenceTransformer


class HelpCenterAgent:
    CHUNKING_VERSION = "parent_child_semantic_v11"
    EMBEDDING_MODEL_NAME = "all-MiniLM-L6-v2"
    RERANKER_MODEL_NAME = "cross-encoder/ms-marco-MiniLM-L-6-v2"
    SEMANTIC_SIMILARITY_THRESHOLD = 0.82
    MAX_SEMANTIC_CHUNK_CHARS = 1400
    MIN_SEMANTIC_WINDOW_CHARS = 320
    TARGET_SEMANTIC_WINDOW_CHARS = 650
    MAX_SEMANTIC_WINDOW_CHARS = 980
    MAX_SEMANTIC_WINDOW_BLOCKS = 3
    LONG_SENTENCE_SPLIT_CHUNK_SIZE = 260
    LONG_SENTENCE_SPLIT_CHUNK_OVERLAP = 40
    CHILD_CHUNK_SIZE = 240
    CHILD_CHUNK_OVERLAP = 40
    RETRIEVAL_K = 12
    PARENT_RERANK_K = 10
    CHILD_FETCH_K = 18
    BM25_PARENT_K = 12
    FILTER_CONFIDENCE_THRESHOLD = 2.0
    MAX_QUERY_VARIANTS = 8
    FINAL_CONTEXT_K_FOCUSED = 4
    FINAL_CONTEXT_K_BROAD = 8
    MIN_BODY_WORDS = 7
    MIN_BODY_CHARS = 45
    TOC_PAGE_SCAN_LIMIT = 4
    MIN_TOPIC_FLAG_SCORE = 3
    DEFAULT_QUERY_REWRITES = {
        "illeness": "illness",
        "holidats": "holidays",
        "how may": "how many",
        "long term": "long-term",
    }
    TOPIC_ALIAS_MAP = {
        "leave_timeoff": "leave_medical",
        "wfh": "remote_work",
        "holiday_calendar": "holiday_calendar",
        "dress_code": "dress_code",
        "visitor_access": "visitor_access",
        "device_usage": "device_usage",
        "data_security": "data_security",
        "termination": "termination",
        "resignation": "resignation",
        "referral_program": "referral_program",
    }
    DEFAULT_INTENT_HINTS = {
        "leave_timeoff": [
            "paid time off",
            "pto",
            "holiday",
            "holidays",
            "sick leave",
            "bereavement leave",
            "jury duty",
            "voting",
            "parental leave",
            "maternity leave",
            "paternity leave",
            "long-term illness",
            "fmla",
        ],
        "wfh": ["work from home", "remote working", "telecommute"],
        "hours_overtime": ["working hours", "overtime", "payroll", "compensation"],
        "benefits": ["benefits", "perks", "health", "insurance"],
        "termination": [
            "termination",
            "discipline",
            "misconduct",
            "dismissal",
            "layoff",
            "without cause",
        ],
        "resignation": ["resignation", "resign", "quit", "notice period", "last day"],
        "device_usage": [
            "company-issued equipment",
            "equipment",
            "laptop",
            "computer",
            "internet",
            "email",
            "cell phone",
            "device",
            "personal use",
        ],
        "visitor_access": [
            "visitor",
            "visitors",
            "sign in",
            "reception",
            "front-office",
            "passes",
            "office access",
        ],
        "holiday_calendar": ["holiday", "holidays", "floating day", "holiday pay"],
        "dress_code": ["dress code", "dress", "appearance", "attire", "grooming"],
        "data_security": [
            "data protection",
            "confidentiality",
            "security",
            "confidential",
            "encryption",
            "privacy",
            "cyber",
        ],
        "referral_program": ["referral", "bonus", "reward", "candidate", "hired"],
    }
    DEFAULT_POLICY_TOPIC_HINTS = {
        "device_usage": [
            "company-issued equipment",
            "equipment",
            "laptop",
            "computer",
            "device",
            "email",
            "internet",
            "cell phone",
            "personal use",
            "company email",
        ],
        "visitor_access": [
            "visitor",
            "visitors",
            "reception",
            "front-office",
            "passes",
            "sign in",
            "show identification",
        ],
        "holiday_calendar": [
            "holiday",
            "holidays",
            "floating day",
            "holiday pay",
            "off-days",
            "observe that holiday",
        ],
        "dress_code": [
            "dress code",
            "business casual",
            "smart casual",
            "casual",
            "clothes",
            "accessories",
            "grooming",
        ],
        "data_security": [
            "data",
            "confidentiality",
            "security",
            "confidential",
            "sensitive data",
            "encryption",
            "backups",
            "access authorization",
        ],
        "leave_medical": [
            "sick leave",
            "long-term illness",
            "fmla",
            "medical",
            "jury duty",
            "bereavement leave",
            "parental leave",
            "pto",
            "time off",
        ],
        "remote_work": [
            "work from home",
            "wfh",
            "remote working",
            "remote work",
            "telecommute",
            "teleworking",
        ],
        "termination": [
            "termination",
            "dismissal",
            "wrongful dismissal",
            "misconduct",
            "for cause",
            "without cause",
            "discipline",
            "layoff",
        ],
        "resignation": [
            "resignation",
            "resign",
            "quit",
            "notice",
            "last day",
            "acceptance of resignation",
        ],
        "referral_program": [
            "referral",
            "bonus",
            "reward",
            "candidate",
            "hard-to-fill role",
            "taxation",
        ],
    }
    TOPIC_PREFERRED_SECTION_MARKERS = {
        "device_usage": [
            "company-issued equipment",
            "cyber security and digital devices",
            "internet usage",
            "corporate email",
            "cell phone",
            "digital devices",
        ],
        "visitor_access": ["workplace visitors", "visitors", "visitor"],
        "holiday_calendar": ["holiday", "holidays"],
        "dress_code": ["dress code", "appearance"],
        "data_security": ["cyber security", "confidentiality", "confidential", "data protection", "security"],
        "leave_medical": [
            "sick leave",
            "long-term illness",
            "paid time off",
            "bereavement leave",
            "jury duty",
            "voting",
            "paternity and maternity leave",
            "parental leave",
        ],
        "remote_work": ["work from home", "remote working"],
        "termination": ["termination", "progressive discipline"],
        "resignation": ["resignation"],
        "referral_program": ["employee referral", "referral"],
    }
    TOPIC_AVOID_SECTION_MARKERS = {
        "device_usage": ["work from home", "workplace visitors", "emergency management", "attendance"],
        "visitor_access": ["drug-free workplace", "dress code", "attendance"],
        "holiday_calendar": ["paternity and maternity leave", "sick leave", "jury duty", "bereavement leave", "paid time off"],
        "data_security": ["social media", "protect our company image", "image and reputation"],
        "termination": ["cobra", "group health benefits", "benefits continuation", "leaving our company"],
        "remote_work": ["benefits and perks", "attendance"],
    }
    DEFAULT_TOPIC_SIGNAL_HINTS = {
        "wfh": {
            "primary_location": [
                "work from home",
                "wfh",
                "telecommute",
                "teleworking",
                "home office",
            ],
            "secondary_location": [
                "remote work",
                "remote working",
                "work remotely",
                "remote arrangement",
            ],
            "allowance": [
                "occasionally",
                "allow",
                "allowed",
                "maximum",
                "per week",
                "per month",
                "days per week",
                "days in advance",
                "recurring",
            ],
            "approval": [
                "manager",
                "supervisor",
                "approval",
                "approve",
                "permission",
                "request",
                "inform",
                "submit",
                "talk to",
            ],
            "notice": [
                "in advance",
                "advance",
                "prior approval",
                "prior notice",
                "notice",
                "submit",
                "request",
            ],
            "conditions": [
                "secure",
                "internet connection",
                "devices",
                "check in",
                "team",
                "collaboration",
                "distractions",
                "time-zone",
            ],
            "adjacent_noise": [
                "benefits",
                "perks",
                "attendance",
                "social media",
                "cell phone",
                "solicitation",
            ],
        },
        "termination": {
            "primary": [
                "termination",
                "terminated",
                "dismissed",
                "dismissal",
                "discharged",
                "layoff",
                "laid off",
                "involuntary separation",
            ],
            "discipline": [
                "discipline",
                "disciplinary",
                "warning",
                "corrective action",
                "performance improvement",
            ],
            "cause": [
                "cause",
                "misconduct",
                "gross misconduct",
                "serious offense",
                "serious offenses",
                "policy violation",
                "without cause",
            ],
            "resignation": [
                "resignation",
                "resign",
                "quit",
                "notice",
                "voluntary separation",
                "last day",
                "forced resignation",
            ],
            "admin": [
                "cobra",
                "reimbursement",
                "benefit continuation",
                "final paycheck",
            ],
        },
        "resignation": {
            "primary": [
                "resignation",
                "resign",
                "quit",
                "notice",
                "voluntary separation",
                "last day",
            ],
            "secondary": [
                "handover",
                "transition",
                "final paycheck",
                "benefits",
            ],
            "termination_noise": [
                "termination",
                "misconduct",
                "discipline",
                "laid off",
                "dismissed",
            ],
        },
    }

    QA_PROMPT = PromptTemplate(
        input_variables=["context", "question"],
        template=(
            "You are an HR policy assistant. Answer using the provided context.\n"
            "If the context includes relevant policy details, provide the best possible answer from it.\n"
            "If the context has template/conditional wording (for example, '[Insert this section ...]'), "
            "state the condition clearly instead of refusing to answer.\n"
            "Use no external knowledge and do not guess beyond the context.\n"
            "Stay close to the retrieved wording and avoid broad interpretations that are not explicitly stated.\n"
            "Do not infer general company rules from optional or conditional sections without naming the condition.\n"
            "If the question asks for a count and the context lists the items, count only the listed items.\n"
            "Use bullet points when listing multiple policy items, conditions, or steps.\n"
            "Keep answers concise and structured.\n"
            "Only if there is no relevant information at all, reply exactly: "
            "'I don't have specific information about that in our airline documents.'\n\n"
            "Context:\n{context}\n\n"
            "Question: {question}\n"
            "Answer:"
        ),
    )

    def __init__(self):
        doc_path = os.environ.get("DOCX_DOC_PATH")
        persist_dir = os.environ.get("PERSIST_DIRECTORY")
        if not doc_path:
            raise ValueError("DOCX_DOC_PATH environment variable is required")
        if not persist_dir:
            raise ValueError("PERSIST_DIRECTORY environment variable is required")

        base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))

        def resolve_path(value: str) -> str:
            if os.path.isabs(value):
                return os.path.abspath(value)
            return os.path.abspath(os.path.join(base_dir, value))

        self._doc_path = resolve_path(doc_path)
        self._persist_dir = os.path.join(
            resolve_path(persist_dir), self.CHUNKING_VERSION
        )
        self._semantic_model = None
        self._reranker = None
        self._parent_docs: Dict[str, Document] = {}
        self._bm25_parent_index: Dict[str, Dict[str, object]] = {}
        self._bm25_avgdl = 0.0
        self._retrieval_k = self._env_int("RETRIEVAL_K", self.RETRIEVAL_K)
        self._final_k_focused = self._env_int(
            "FINAL_CONTEXT_K_FOCUSED", self.FINAL_CONTEXT_K_FOCUSED
        )
        self._final_k_broad = self._env_int(
            "FINAL_CONTEXT_K_BROAD", self.FINAL_CONTEXT_K_BROAD
        )
        self._toc_page_scan_limit = self._env_int(
            "TOC_PAGE_SCAN_LIMIT", self.TOC_PAGE_SCAN_LIMIT
        )

        if not os.path.isdir(self._doc_path):
            raise ValueError(
                f"DOCX_DOC_PATH does not exist or is not a directory: {self._doc_path}"
            )

        os.makedirs(self._persist_dir, exist_ok=True)
        self._docx_db = self._create_index(directory=self._doc_path)
        self._qa_chain = self._create_qa_chain()

    def _create_index(self, directory: str):
        embeddings = SentenceTransformerEmbeddings(model_name=self.EMBEDDING_MODEL_NAME)
        if os.path.exists(self._persist_dir) and os.listdir(self._persist_dir):
            print(f"Loading existing index from {self._persist_dir}")
            existing_db = Chroma(
                persist_directory=self._persist_dir,
                embedding_function=embeddings,
            )
            collection = getattr(existing_db, "_collection", None)
            self._parent_docs = self._load_parent_docs()
            self._build_bm25_parent_index()
            if collection is not None and collection.count() > 0 and self._parent_docs:
                return existing_db
            print(f"Existing index at {self._persist_dir} is empty. Rebuilding it.")

        print(f"Building new index from {directory}")
        child_docs, parent_docs = self.split_docs(self.load_docs(directory))
        if not child_docs or not parent_docs:
            raise ValueError(f"No documents were loaded from {directory}")
        self._parent_docs = {
            str((doc.metadata or {}).get("parent_doc_id")): doc for doc in parent_docs
        }
        self._save_parent_docs(parent_docs)
        self._build_bm25_parent_index()
        vectordb = Chroma.from_documents(
            documents=child_docs,
            embedding=embeddings,
            persist_directory=self._persist_dir,
        )
        vectordb.persist()
        return vectordb

    def docx_retriever(self):
        return self._docx_db.as_retriever()

    def _parent_store_path(self) -> str:
        return os.path.join(self._persist_dir, "parent_docs.json")

    @classmethod
    def _serialize_doc(cls, doc: Document) -> Dict[str, object]:
        return {
            "page_content": doc.page_content,
            "metadata": dict(doc.metadata or {}),
        }

    @classmethod
    def _deserialize_doc(cls, payload: Dict[str, object]) -> Document:
        return Document(
            page_content=str(payload.get("page_content", "")),
            metadata=dict(payload.get("metadata") or {}),
        )

    @staticmethod
    def _file_key(source_path: str) -> str:
        """Stable short key for a source file, used in per-file parent IDs."""
        basename = os.path.basename(source_path or "")
        return hashlib.sha1(basename.encode("utf-8")).hexdigest()[:8]

    def _save_parent_docs(self, docs: Sequence[Document]) -> None:
        # Atomic on Windows: write to a temp file, then os.replace.
        path = self._parent_store_path()
        tmp_path = path + ".tmp"
        with open(tmp_path, "w", encoding="utf-8") as handle:
            json.dump([self._serialize_doc(doc) for doc in docs], handle, ensure_ascii=False)
        os.replace(tmp_path, path)

    def _persist_parent_state(self) -> None:
        """Save the current in-memory parent dict and refresh BM25."""
        self._save_parent_docs(list(self._parent_docs.values()))
        self._build_bm25_parent_index()

    # ── Incremental indexing ──────────────────────────────────────────
    def add_file(self, file_path: str) -> int:
        """Index a single new file without rebuilding the whole KB.

        Returns the number of child chunks added. Raises on any failure —
        the caller is expected to fall back to a full rebuild.
        """
        file_path = os.path.abspath(file_path)
        documents = self._load_single_file(file_path)
        if not documents:
            raise ValueError(f"No content extracted from {file_path}")
        child_docs, parent_docs = self.split_docs(documents)
        if not child_docs or not parent_docs:
            raise ValueError(f"No indexable chunks produced from {file_path}")

        self._docx_db.add_documents(child_docs)
        self._docx_db.persist()

        for doc in parent_docs:
            parent_id = str((doc.metadata or {}).get("parent_doc_id", "")).strip()
            if parent_id:
                self._parent_docs[parent_id] = doc
        self._persist_parent_state()
        return len(child_docs)

    def remove_file(self, file_path: str) -> int:
        """Remove all chunks belonging to one source file from the index.

        Returns the number of parent chunks removed. Raises on failure —
        the caller is expected to fall back to a full rebuild.
        """
        file_path = os.path.abspath(file_path)
        collection = getattr(self._docx_db, "_collection", None)
        if collection is None:
            raise RuntimeError("Chroma collection unavailable for delete")

        # Children store the same absolute path in "source" metadata.
        collection.delete(where={"source": {"$eq": file_path}})
        self._docx_db.persist()

        removed_ids = [
            parent_id
            for parent_id, doc in self._parent_docs.items()
            if os.path.abspath(str((doc.metadata or {}).get("source", ""))) == file_path
        ]
        for parent_id in removed_ids:
            del self._parent_docs[parent_id]
        self._persist_parent_state()
        return len(removed_ids)

    def _load_single_file(self, file_path: str) -> List[Document]:
        """Run the appropriate loader for one file (mirrors load_docs routing)."""
        lower = os.path.basename(file_path).lower()
        ext = os.path.splitext(lower)[1]
        if lower.endswith(".docx"):
            return self._load_docx(file_path)
        if lower.endswith(".pdf"):
            return self._load_pdf(file_path)
        if ext in self.GENERIC_TEXT_EXTENSIONS:
            return self._load_text(file_path)
        if ext in (".pptx", ".xlsx"):
            return self._load_office(file_path)
        raise ValueError(f"Unsupported file type for incremental indexing: {ext}")

    def _load_parent_docs(self) -> Dict[str, Document]:
        path = self._parent_store_path()
        if not os.path.exists(path):
            return {}
        try:
            with open(path, "r", encoding="utf-8") as handle:
                payload = json.load(handle)
        except Exception:
            return {}
        parent_docs: Dict[str, Document] = {}
        if isinstance(payload, list):
            for item in payload:
                if not isinstance(item, dict):
                    continue
                doc = self._deserialize_doc(item)
                parent_id = str((doc.metadata or {}).get("parent_doc_id", "")).strip()
                if parent_id:
                    parent_docs[parent_id] = doc
        return parent_docs

    def _get_reranker(self) -> CrossEncoder:
        if self._reranker is None:
            self._reranker = CrossEncoder(self.RERANKER_MODEL_NAME)
        return self._reranker

    @classmethod
    def _bm25_tokens(cls, text: str) -> List[str]:
        return [
            token
            for token in re.findall(r"[a-zA-Z0-9]+", cls._normalize_text(text).lower())
            if len(token) > 1
        ]

    def _build_bm25_parent_index(self) -> None:
        self._bm25_parent_index = {}
        lengths: List[int] = []
        for parent_id, doc in self._parent_docs.items():
            metadata = doc.metadata or {}
            lexical_text = " ".join(
                [
                    str(metadata.get("parent_heading", "")),
                    str(metadata.get("section_path", "")),
                    str(metadata.get("heading_context", "")),
                    self._source_content(doc),
                ]
            )
            tokens = self._bm25_tokens(lexical_text)
            token_counts = Counter(tokens)
            doc_len = sum(token_counts.values())
            lengths.append(doc_len)
            self._bm25_parent_index[parent_id] = {
                "doc": doc,
                "counts": token_counts,
                "len": doc_len,
            }
        self._bm25_avgdl = (sum(lengths) / len(lengths)) if lengths else 0.0

    def _bm25_parent_results(
        self, query: str, query_topics: Sequence[str]
    ) -> List[Tuple[float, Document, float]]:
        if not self._bm25_parent_index:
            return []

        query_terms = self._bm25_tokens(query)
        if not query_terms:
            return []

        doc_freq: Dict[str, int] = {}
        for term in set(query_terms):
            doc_freq[term] = sum(
                1
                for item in self._bm25_parent_index.values()
                if item["counts"].get(term, 0) > 0
            )

        total_docs = max(1, len(self._bm25_parent_index))
        k1 = 1.5
        b = 0.75
        scored: List[Tuple[float, Document, float]] = []
        primary_topic = query_topics[0] if query_topics else ""

        for item in self._bm25_parent_index.values():
            doc = item["doc"]
            metadata = doc.metadata or {}
            counts: Counter = item["counts"]
            doc_len = int(item["len"])
            score = 0.0
            for term in query_terms:
                freq = counts.get(term, 0)
                if freq <= 0:
                    continue
                df = doc_freq.get(term, 0)
                idf = math.log(1 + ((total_docs - df + 0.5) / (df + 0.5)))
                denom = freq + k1 * (
                    1 - b + b * (doc_len / max(self._bm25_avgdl, 1.0))
                )
                score += idf * ((freq * (k1 + 1)) / max(denom, 1e-9))

            if score <= 0:
                continue

            section_family_text = self._section_family_text(doc)
            preferred_hits = self._count_markers(
                section_family_text, self.TOPIC_PREFERRED_SECTION_MARKERS.get(primary_topic, [])
            )
            avoid_hits = self._count_markers(
                section_family_text, self.TOPIC_AVOID_SECTION_MARKERS.get(primary_topic, [])
            )
            topic_bonus = 0.4 * self._doc_topic_overlap(doc, query_topics)
            focus_bonus = (
                0.5
                if primary_topic
                and str(metadata.get("policy_focus", "")) == primary_topic
                else 0.0
            )
            adjusted = score + topic_bonus + focus_bonus + (0.35 * preferred_hits) - (0.5 * avoid_hits)
            scored.append((-adjusted, doc, -adjusted))

        scored.sort(key=lambda item: item[0])
        return scored[: self.BM25_PARENT_K]

    def _run_query(self, query: str):
        import time as _time

        normalized_query = self._normalize_query(query)
        _t0 = _time.perf_counter()
        ranked_results = self._retrieve_candidates(normalized_query)
        matching_docs, sources = self._select_final_context(normalized_query, ranked_results)
        retrieval_ms = int((_time.perf_counter() - _t0) * 1000)

        # Plain-data trace details for observability (consumed by api.py);
        # no tracing SDK coupling in this layer.
        trace = {
            "retrieval_ms": retrieval_ms,
            "generation_ms": 0,
            "candidate_count": len(ranked_results or []),
            "answered_by": "no_context",
        }

        if not matching_docs:
            return {
                "answer": "I don't have specific information about that in our airline documents.",
                "sources": [],
                "trace": trace,
            }

        direct_answer = self._answer_from_context(query, normalized_query, matching_docs)
        if direct_answer is not None:
            trace["answered_by"] = "direct"
            return {"answer": direct_answer, "sources": sources, "trace": trace}

        _t1 = _time.perf_counter()
        answer = self._qa_chain.run(
            input_documents=matching_docs, question=normalized_query
        )
        trace["generation_ms"] = int((_time.perf_counter() - _t1) * 1000)
        trace["answered_by"] = "llm"
        return {"answer": answer, "sources": sources, "trace": trace}

    def _select_final_context(self, query: str, ranked_results):
        matching_docs = []
        sources = []
        seen_content = set()
        seen_sections = {}
        final_k = self._final_k_for_query(query)
        query_topics = self._query_topics(query)
        primary_topic = query_topics[0] if query_topics else ""
        has_topic_matches = any(
            self._doc_topic_overlap(doc, query_topics) > 0 for _, doc, _ in ranked_results
        )

        def try_add(doc, score, allow_duplicate_section: bool = False) -> bool:
            if self._is_low_value_candidate(doc):
                return False
            if self._should_skip_topic_candidate(query, doc, query_topics, has_topic_matches):
                return False
            normalized = self._normalize_text(doc.page_content).lower()
            if normalized in seen_content:
                return False

            # Avoid returning 3-4 chunks from the same section/heading.
            meta = doc.metadata or {}
            section_key = (
                self._normalize_text(meta.get("parent_heading", "")) or self._normalize_text(meta.get("section_path", ""))
            ).lower()
            if section_key and not allow_duplicate_section and seen_sections.get(section_key, 0) >= 1:
                return False

            seen_content.add(normalized)
            if section_key:
                seen_sections[section_key] = seen_sections.get(section_key, 0) + 1
            matching_docs.append(doc)
            sources.append(
                {
                    "content": self._source_content(doc),
                    "metadata": doc.metadata,
                    "score": score,
                }
            )
            return True

        preferred_candidates = self._preferred_context_candidates(
            query, ranked_results, query_topics, primary_topic
        )
        for _, doc, score in preferred_candidates:
            if try_add(doc, score, allow_duplicate_section=True) and len(matching_docs) >= min(final_k, 2):
                break

        # First prefer unique sections so broad queries span multiple policy areas.
        for _, doc, score in ranked_results:
            if try_add(doc, score) and len(matching_docs) >= final_k:
                return matching_docs, sources

        # If we still need more context, allow a second chunk from the same section.
        for _, doc, score in ranked_results:
            if try_add(doc, score, allow_duplicate_section=True) and len(matching_docs) >= final_k:
                break

        return matching_docs, sources

    def _preferred_context_candidates(
        self,
        query: str,
        ranked_results,
        query_topics: Sequence[str],
        primary_topic: str,
    ):
        if not primary_topic:
            return []

        preferred = []
        require_topical = any(
            self._doc_topic_overlap(doc, query_topics) > 0 for _, doc, _ in ranked_results
        )
        for item in ranked_results:
            _, doc, _ = item
            metadata = doc.metadata or {}
            topic_overlap = self._doc_topic_overlap(doc, query_topics)
            policy_focus = str(metadata.get("policy_focus", ""))
            section_text = self._section_family_text(doc)
            preferred_hits = self._count_markers(
                section_text, self.TOPIC_PREFERRED_SECTION_MARKERS.get(primary_topic, [])
            )
            avoid_hits = self._count_markers(
                section_text, self.TOPIC_AVOID_SECTION_MARKERS.get(primary_topic, [])
            )
            if require_topical and topic_overlap == 0:
                continue
            if policy_focus == primary_topic or preferred_hits > 0:
                preferred.append(
                    (
                        item[0] - min(0.18, 0.06 * preferred_hits) + min(0.18, 0.08 * avoid_hits),
                        doc,
                        item[2],
                    )
                )
        preferred.sort(key=lambda item: item[0])
        return preferred

    @classmethod
    def _section_key_for_doc(cls, doc: Document) -> str:
        metadata = doc.metadata or {}
        return (
            cls._normalize_text(metadata.get("parent_heading", ""))
            or cls._normalize_text(metadata.get("section_path", ""))
        ).lower()

    @classmethod
    def _section_context_doc(
        cls, docs: Sequence[Document], section_score: float
    ) -> Tuple[float, Document, float] | None:
        if not docs:
            return None
        metadata = dict((docs[0].metadata or {}))
        raw_parts = [cls._normalize_text(cls._source_content(doc)) for doc in docs]
        raw_parts = [part for part in raw_parts if part]
        if not raw_parts:
            return None

        raw_content = cls._normalize_text(" ".join(raw_parts))
        if not cls._has_substantive_body(raw_content, metadata):
            return None

        section_doc = Document(
            page_content=cls._semantic_chunk_text(raw_content, metadata),
            metadata=cls._copy_metadata(
                metadata,
                raw_content=raw_content,
                chunk_strategy="section_context",
                chunk_stage="section_context",
                assembled_from=len(docs),
            ),
        )
        return (section_score - 0.01, section_doc, section_score - 0.01)

    @classmethod
    def _section_fallback_candidates(cls, query: str, ranked_results):
        section_groups: Dict[str, List[Tuple[float, Document, float]]] = {}
        section_order: List[str] = []
        top_slice = ranked_results[: max(8, cls.FINAL_CONTEXT_K_BROAD)]

        for item in top_slice:
            _, doc, _ = item
            if cls._is_low_value_candidate(doc) or cls._should_skip_topic_candidate(query, doc):
                continue
            section_key = cls._section_key_for_doc(doc)
            if not section_key:
                continue
            if section_key not in section_groups:
                section_groups[section_key] = []
                section_order.append(section_key)
            section_groups[section_key].append(item)

        fallbacks: List[Tuple[float, Document, float]] = []
        for section_key in section_order:
            items = section_groups[section_key]
            if len(items) < 2:
                continue

            docs = [items[0][1]]
            first_raw = cls._source_content(items[0][1])
            for item in items[1:]:
                candidate_doc = item[1]
                prev_meta = docs[-1].metadata or {}
                cand_meta = candidate_doc.metadata or {}
                adjacent = (
                    prev_meta.get("section_path") == cand_meta.get("section_path")
                    and prev_meta.get("page_number") == cand_meta.get("page_number")
                    and cand_meta.get("section_window_index")
                    == prev_meta.get("section_window_index", 0) + 1
                )
                if adjacent:
                    docs.append(candidate_doc)
                if len(docs) >= 2:
                    break

            if len(docs) < 2:
                continue

            combined_raw = cls._normalize_text(" ".join(cls._source_content(doc) for doc in docs))
            first_terms = cls._query_terms(query) + cls._intent_hint_terms(cls._topic_for_query(query))
            first_hits = cls._count_term_hits(first_raw, first_terms)
            combined_hits = cls._count_term_hits(combined_raw, first_terms)
            if combined_hits <= first_hits and len(first_raw) >= cls.TARGET_SEMANTIC_WINDOW_CHARS:
                continue

            section_candidate = cls._section_context_doc(docs, items[0][0])
            if section_candidate is not None:
                fallbacks.append(section_candidate)

        fallbacks.sort(key=lambda item: item[0])
        return fallbacks

    @classmethod
    def _should_skip_topic_candidate(
        cls,
        query: str,
        doc: Document,
        query_topics: Sequence[str] | None = None,
        has_topic_matches: bool = False,
    ) -> bool:
        if not query_topics:
            return False
        if not has_topic_matches:
            return False
        overlap = cls._doc_topic_overlap(doc, query_topics)
        if overlap == 0:
            return True
        primary_topic = query_topics[0]
        if cls._count_markers(
            cls._section_family_text(doc),
            cls.TOPIC_AVOID_SECTION_MARKERS.get(primary_topic, []),
        ) > 0 and cls._count_markers(
            cls._section_family_text(doc),
            cls.TOPIC_PREFERRED_SECTION_MARKERS.get(primary_topic, []),
        ) == 0:
            return True
        return False

    def _retrieve_candidates(self, query: str):
        query_topics = self._query_topics(query)
        hint_terms: List[str] = []
        for topic in query_topics:
            hint_terms.extend(self._intent_hint_terms(self._legacy_topic_name(topic)))
        queries = self._build_query_variants(query, hint_terms)

        child_results = self._child_retrieval_results(queries, query_topics)
        dense_parent_candidates = self._collapse_child_results_to_parents(query, child_results, query_topics)
        sparse_parent_candidates = self._bm25_parent_results(query, query_topics)
        fused_parent_candidates = self._fuse_parent_candidates(
            dense_parent_candidates, sparse_parent_candidates, query_topics
        )
        return self._rerank_parent_candidates(query, fused_parent_candidates, query_topics)

    def _fuse_parent_candidates(
        self,
        dense_parent_candidates: Sequence[Tuple[float, Document, float]],
        sparse_parent_candidates: Sequence[Tuple[float, Document, float]],
        query_topics: Sequence[str],
    ) -> List[Tuple[float, Document, float]]:
        fused: Dict[str, Dict[str, object]] = {}
        primary_topic = query_topics[0] if query_topics else ""

        for rank, (score, doc, raw_score) in enumerate(dense_parent_candidates, start=1):
            parent_id = str((doc.metadata or {}).get("parent_doc_id", "")).strip()
            if not parent_id:
                continue
            fused[parent_id] = {
                "doc": doc,
                "dense_score": float(score),
                "raw_score": float(raw_score),
                "rrf": 1.0 / (10 + rank),
                "dense_rank": rank,
                "sparse_rank": None,
            }

        for rank, (score, doc, raw_score) in enumerate(sparse_parent_candidates, start=1):
            parent_id = str((doc.metadata or {}).get("parent_doc_id", "")).strip()
            if not parent_id:
                continue
            if parent_id not in fused:
                fused[parent_id] = {
                    "doc": doc,
                    "dense_score": 1.4,
                    "raw_score": float(raw_score),
                    "rrf": 0.0,
                    "dense_rank": None,
                    "sparse_rank": rank,
                }
            fused[parent_id]["rrf"] = float(fused[parent_id]["rrf"]) + (1.0 / (10 + rank))
            fused[parent_id]["sparse_rank"] = rank
            fused[parent_id]["raw_score"] = min(float(fused[parent_id]["raw_score"]), float(raw_score))

        results: List[Tuple[float, Document, float]] = []
        for item in fused.values():
            doc = item["doc"]
            metadata = doc.metadata or {}
            section_family_text = self._section_family_text(doc)
            preferred_hits = self._count_markers(
                section_family_text, self.TOPIC_PREFERRED_SECTION_MARKERS.get(primary_topic, [])
            )
            avoid_hits = self._count_markers(
                section_family_text, self.TOPIC_AVOID_SECTION_MARKERS.get(primary_topic, [])
            )
            focus_bonus = (
                0.08
                if primary_topic and str(metadata.get("policy_focus", "")) == primary_topic
                else 0.0
            )
            fused_score = (
                float(item["dense_score"])
                - (1.4 * float(item["rrf"]))
                - min(0.18, 0.06 * preferred_hits)
                - focus_bonus
                + min(0.22, 0.08 * avoid_hits)
            )
            results.append((fused_score, doc, float(item["raw_score"])))

        results.sort(key=lambda item: item[0])
        return results[: max(self.PARENT_RERANK_K * 2, self._retrieval_k)]

    @staticmethod
    def _env_int(name: str, default: int) -> int:
        raw = os.environ.get(name)
        if raw is None or str(raw).strip() == "":
            return int(default)
        try:
            return int(str(raw).strip())
        except ValueError:
            return int(default)

    @classmethod
    def _legacy_topic_name(cls, topic: str) -> str:
        reverse = {value: key for key, value in cls.TOPIC_ALIAS_MAP.items()}
        return reverse.get(topic, topic)

    @classmethod
    def _query_topic_scores(cls, query: str) -> List[Tuple[str, float]]:
        normalized = cls._normalize_query(query)
        topic_scores: List[Tuple[str, float]] = []
        seen_topics = sorted(set(cls.TOPIC_ALIAS_MAP.values()))

        for topic in seen_topics:
            markers: List[str] = []
            markers.extend(cls.DEFAULT_POLICY_TOPIC_HINTS.get(topic, []))
            legacy = cls._legacy_topic_name(topic)
            markers.extend(cls._intent_hint_terms(legacy))
            score = float(cls._count_markers(normalized, markers))
            if topic == "holiday_calendar" and cls._is_holiday_count_query(normalized):
                score += 1.0
            if topic == "leave_medical" and cls._is_timeoff_leave_query(normalized):
                score += 1.0
            if topic == "remote_work" and cls._is_wfh_query(normalized):
                score += 1.0
            if topic == "termination" and cls._is_termination_query(normalized):
                score += 1.0
            if topic == "resignation" and cls._is_resignation_query(normalized):
                score += 1.0
            if score > 0:
                topic_scores.append((topic, score))

        topic_scores.sort(key=lambda item: item[1], reverse=True)
        return topic_scores

    @classmethod
    def _query_topics(cls, query: str) -> List[str]:
        scores = cls._query_topic_scores(query)
        if not scores:
            return []
        top_topic, top_score = scores[0]
        selected = [top_topic]
        for topic, score in scores[1:3]:
            if score >= 2.0 and score >= max(1.5, top_score * 0.75):
                selected.append(topic)
        return selected

    @classmethod
    def _query_topic_filters(cls, query: str) -> List[Tuple[str, Dict[str, str]]]:
        filters: List[Tuple[str, Dict[str, str]]] = []
        scores = cls._query_topic_scores(query)
        if cls._is_holiday_count_query(query):
            scores = [item for item in scores if item[0] == "holiday_calendar"]
        for idx, (topic, score) in enumerate(scores):
            if score < cls.FILTER_CONFIDENCE_THRESHOLD:
                continue
            if idx > 0 and score < scores[0][1]:
                break
            filters.append((topic, {f"topic_{topic}": "1"}))
            if len(filters) >= 2:
                break
        return filters

    @classmethod
    def _doc_topic_overlap(cls, doc: Document, query_topics: Sequence[str]) -> int:
        metadata = doc.metadata or {}
        return sum(1 for topic in query_topics if str(metadata.get(f"topic_{topic}", "0")) == "1")

    @classmethod
    def _section_family_text(cls, doc: Document) -> str:
        metadata = doc.metadata or {}
        return cls._normalize_text(
            " ".join(
                [
                    str(metadata.get("parent_heading", "")),
                    str(metadata.get("section_path", "")),
                    str(metadata.get("heading_context", "")),
                    str(metadata.get("policy_focus", "")),
                ]
            )
        ).lower()

    def _child_retrieval_results(self, queries: Sequence[str], query_topics: Sequence[str]):
        results: List[Tuple[Document, float]] = []
        seen = set()
        query_filters = self._query_topic_filters(" ".join(queries))
        k_each = max(4, int(self._retrieval_k))
        combined_query = self._normalize_text(" ".join(queries))
        combined_terms = self._query_terms(combined_query)

        def add_result(doc: Document, score: float) -> None:
            metadata = doc.metadata or {}
            key = (
                str(metadata.get("parent_doc_id", "")),
                str(metadata.get("child_index", "")),
                self._normalize_text(self._source_content(doc)).lower(),
            )
            if key in seen:
                return
            lexical_hits = self._count_term_hits(
                " ".join(
                    [
                        str(metadata.get("parent_heading", "")),
                        str(metadata.get("section_path", "")),
                        self._source_content(doc),
                    ]
                ),
                combined_terms,
            )
            topic_bonus = 0.03 * self._doc_topic_overlap(doc, query_topics)
            seen.add(key)
            results.append((doc, float(score) - min(0.12, 0.02 * lexical_hits) - topic_bonus))

        for q in queries:
            try:
                for doc, score in self._docx_db.similarity_search_with_score(q, k=k_each):
                    add_result(doc, score)
            except Exception:
                pass

            try:
                mmr_docs = self._docx_db.max_marginal_relevance_search(
                    q,
                    k=max(4, min(k_each, 8)),
                    fetch_k=max(self.CHILD_FETCH_K, k_each * 2),
                )
                for idx, doc in enumerate(mmr_docs):
                    add_result(doc, 0.78 + (0.02 * idx))
            except Exception:
                pass

            for _, topic_filter in query_filters:
                try:
                    for doc, score in self._docx_db.similarity_search_with_score(
                        q, k=k_each, filter=topic_filter
                    ):
                        add_result(doc, score - 0.02)
                except Exception:
                    pass
                try:
                    mmr_docs = self._docx_db.max_marginal_relevance_search(
                        q,
                        k=max(4, min(k_each, 8)),
                        fetch_k=max(self.CHILD_FETCH_K, k_each * 2),
                        filter=topic_filter,
                    )
                    for idx, doc in enumerate(mmr_docs):
                        add_result(doc, 0.68 + (0.02 * idx))
                except Exception:
                    pass

        results.sort(key=lambda item: item[1])
        return results

    def _collapse_child_results_to_parents(
        self,
        query: str,
        child_results: Sequence[Tuple[Document, float]],
        query_topics: Sequence[str],
    ) -> List[Tuple[float, Document, float]]:
        parent_scores: Dict[str, Dict[str, object]] = {}
        query_terms = self._query_terms(query)

        for child_doc, score in child_results:
            metadata = child_doc.metadata or {}
            parent_id = str(metadata.get("parent_doc_id", "")).strip()
            if not parent_id or parent_id not in self._parent_docs:
                continue
            parent_doc = self._parent_docs[parent_id]
            if parent_id not in parent_scores:
                parent_scores[parent_id] = {
                    "doc": parent_doc,
                    "best_score": float(score),
                    "match_count": 0,
                    "topic_overlap": self._doc_topic_overlap(parent_doc, query_topics),
                    "child_indexes": set(),
                }
            record = parent_scores[parent_id]
            record["best_score"] = min(float(record["best_score"]), float(score))
            record["match_count"] = int(record["match_count"]) + 1
            child_index = metadata.get("child_index")
            if child_index is not None:
                record["child_indexes"].add(str(child_index))

        collapsed: List[Tuple[float, Document, float]] = []
        for record in parent_scores.values():
            parent_doc = record["doc"]
            best_score = float(record["best_score"])
            match_count = int(record["match_count"])
            topic_overlap = int(record["topic_overlap"])
            primary_topic = query_topics[0] if query_topics else ""
            heading_text = self._normalize_text(
                " ".join(
                    [
                        str((parent_doc.metadata or {}).get("parent_heading", "")),
                        str((parent_doc.metadata or {}).get("section_path", "")),
                    ]
                )
            )
            body_text = self._source_content(parent_doc)
            heading_hits = self._count_term_hits(heading_text, query_terms)
            body_hits = self._count_term_hits(body_text[:1800], query_terms)
            section_family_text = self._section_family_text(parent_doc)
            preferred_hits = self._count_markers(
                section_family_text, self.TOPIC_PREFERRED_SECTION_MARKERS.get(primary_topic, [])
            )
            avoid_hits = self._count_markers(
                section_family_text, self.TOPIC_AVOID_SECTION_MARKERS.get(primary_topic, [])
            )
            focus_bonus = (
                0.10
                if str((parent_doc.metadata or {}).get("policy_focus", "")) in set(query_topics)
                else 0.0
            )
            focus_penalty = (
                0.12
                if query_topics
                and str((parent_doc.metadata or {}).get("policy_focus", ""))
                and str((parent_doc.metadata or {}).get("policy_focus", "")) not in set(query_topics)
                else 0.0
            )
            diversity_bonus = min(0.10, 0.03 * len(record["child_indexes"]))
            retrieval_score = (
                best_score
                - min(0.18, 0.04 * match_count)
                - min(0.12, 0.05 * topic_overlap)
                - min(0.16, 0.05 * heading_hits)
                - min(0.12, 0.03 * body_hits)
                - focus_bonus
                - min(0.18, 0.06 * preferred_hits)
                - diversity_bonus
                + min(0.22, 0.09 * avoid_hits)
                + focus_penalty
            )
            collapsed.append((retrieval_score, parent_doc, best_score))

        collapsed.sort(key=lambda item: item[0])
        return collapsed[: max(self.PARENT_RERANK_K * 2, self._retrieval_k)]

    def _rerank_parent_candidates(
        self,
        query: str,
        parent_candidates: Sequence[Tuple[float, Document, float]],
        query_topics: Sequence[str],
    ) -> List[Tuple[float, Document, float]]:
        if not parent_candidates:
            return []

        reranker = self._get_reranker()
        query_terms = self._query_terms(query)
        pairs = [
            (query, self._source_content(doc)[:1800])
            for _, doc, _ in parent_candidates[: self.PARENT_RERANK_K]
        ]
        rerank_scores = reranker.predict(pairs) if pairs else []

        ranked: List[Tuple[float, Document, float]] = []
        for idx, (retrieval_score, doc, child_score) in enumerate(parent_candidates[: self.PARENT_RERANK_K]):
            rerank_score = float(rerank_scores[idx]) if idx < len(rerank_scores) else 0.0
            metadata = doc.metadata or {}
            primary_topic = query_topics[0] if query_topics else ""
            heading_text = self._normalize_text(
                " ".join(
                    [
                        str(metadata.get("parent_heading", "")),
                        str(metadata.get("section_path", "")),
                    ]
                )
            )
            body_text = self._source_content(doc)
            heading_hits = self._count_term_hits(heading_text, query_terms)
            body_hits = self._count_term_hits(body_text[:1600], query_terms)
            section_family_text = self._section_family_text(doc)
            preferred_hits = self._count_markers(
                section_family_text, self.TOPIC_PREFERRED_SECTION_MARKERS.get(primary_topic, [])
            )
            avoid_hits = self._count_markers(
                section_family_text, self.TOPIC_AVOID_SECTION_MARKERS.get(primary_topic, [])
            )
            topic_bonus = 0.10 * self._doc_topic_overlap(doc, query_topics)
            focus_bonus = 0.12 if str(metadata.get("policy_focus", "")) in set(query_topics) else 0.0
            focus_penalty = (
                0.12
                if query_topics
                and str(metadata.get("policy_focus", ""))
                and str(metadata.get("policy_focus", "")) not in set(query_topics)
                else 0.0
            )
            final_score = (
                retrieval_score
                - (0.18 * rerank_score)
                - topic_bonus
                - focus_bonus
                - min(0.16, 0.05 * preferred_hits)
                - min(0.14, 0.04 * heading_hits)
                - min(0.10, 0.02 * body_hits)
                + min(0.22, 0.08 * avoid_hits)
                + focus_penalty
            )
            ranked.append((final_score, doc, child_score))

        for retrieval_score, doc, child_score in parent_candidates[self.PARENT_RERANK_K :]:
            metadata = doc.metadata or {}
            primary_topic = query_topics[0] if query_topics else ""
            heading_text = self._normalize_text(
                " ".join(
                    [
                        str(metadata.get("parent_heading", "")),
                        str(metadata.get("section_path", "")),
                    ]
                )
            )
            heading_hits = self._count_term_hits(heading_text, query_terms)
            section_family_text = self._section_family_text(doc)
            preferred_hits = self._count_markers(
                section_family_text, self.TOPIC_PREFERRED_SECTION_MARKERS.get(primary_topic, [])
            )
            avoid_hits = self._count_markers(
                section_family_text, self.TOPIC_AVOID_SECTION_MARKERS.get(primary_topic, [])
            )
            topic_bonus = 0.06 * self._doc_topic_overlap(doc, query_topics)
            focus_bonus = 0.08 if str(metadata.get("policy_focus", "")) in set(query_topics) else 0.0
            focus_penalty = (
                0.10
                if query_topics
                and str(metadata.get("policy_focus", ""))
                and str(metadata.get("policy_focus", "")) not in set(query_topics)
                else 0.0
            )
            ranked.append(
                (
                    retrieval_score
                    - topic_bonus
                    - focus_bonus
                    - min(0.10, 0.04 * preferred_hits)
                    - min(0.08, 0.03 * heading_hits)
                    + min(0.18, 0.07 * avoid_hits)
                    + focus_penalty,
                    doc,
                    child_score,
                )
            )

        ranked.sort(key=lambda item: item[0])
        return ranked

    @classmethod
    def _configured_signal_hints(cls, topic: str) -> Dict[str, List[str]]:
        raw = os.environ.get("TOPIC_SIGNAL_HINTS_JSON")
        if raw:
            try:
                parsed = json.loads(raw)
                if isinstance(parsed, dict):
                    topic_data = parsed.get(topic) or parsed.get(topic.lower()) or {}
                    if isinstance(topic_data, dict):
                        cleaned: Dict[str, List[str]] = {}
                        for key, values in topic_data.items():
                            if isinstance(values, str):
                                values = [values]
                            if isinstance(values, list):
                                cleaned[key] = [
                                    cls._normalize_text(str(value)).lower()
                                    for value in values
                                    if cls._normalize_text(str(value))
                                ]
                        if cleaned:
                            return cleaned
            except Exception:
                pass
        return {
            key: [cls._normalize_text(value).lower() for value in values if value]
            for key, values in cls.DEFAULT_TOPIC_SIGNAL_HINTS.get(topic, {}).items()
        }

    @classmethod
    def _marker_pattern(cls, marker: str) -> str:
        parts = [
            re.escape(part)
            for part in re.split(r"\s+", cls._normalize_text(marker).lower())
            if part
        ]
        if not parts:
            return ""
        core = r"\s+".join(parts)
        return rf"(?<![A-Za-z0-9]){core}(?![A-Za-z0-9])"

    @classmethod
    def _marker_present(cls, text: str, marker: str) -> bool:
        pattern = cls._marker_pattern(marker)
        if not pattern:
            return False
        return re.search(pattern, cls._normalize_text(text).lower()) is not None

    @classmethod
    def _has_any_marker(cls, text: str, markers: Sequence[str]) -> bool:
        return any(cls._marker_present(text, marker) for marker in markers if marker)

    @classmethod
    def _count_markers(cls, text: str, markers: Sequence[str]) -> int:
        return sum(1 for marker in markers if marker and cls._marker_present(text, marker))

    @classmethod
    def _wfh_duration_hits(cls, text: str) -> int:
        normalized = cls._normalize_text(text).lower()
        patterns = [
            r"\b\d+\s+(day|days|week|weeks|month|months)\b",
            r"\b(one|two|three|four|five)\s+(day|days|week|weeks|month|months)\b",
            r"\bper (day|week|month|year)\b",
            r"\bconsecutive\b",
        ]
        return sum(1 for pattern in patterns if re.search(pattern, normalized))

    @classmethod
    def _termination_query_markers(cls) -> List[str]:
        return [
            "terminate",
            "termination",
            "fired",
            "dismiss",
            "laid off",
            "layoff",
            "discharge",
            "involuntary",
        ]

    @classmethod
    def _resignation_query_markers(cls) -> List[str]:
        return ["quit", "resign", "resignation", "notice", "last day", "step down"]

    @classmethod
    def _is_wfh_query(cls, query: str) -> bool:
        return cls._has_any_marker(
            query,
            ["work from home", "wfh", "telecommute", "teleworking", "remote working", "work remotely"],
        )

    @classmethod
    def _is_resignation_query(cls, query: str) -> bool:
        return cls._has_any_marker(query, cls._resignation_query_markers())

    @classmethod
    def _is_termination_query(cls, query: str) -> bool:
        return cls._has_any_marker(query, cls._termination_query_markers())

    @classmethod
    def _is_simple_wfh_rules_query(cls, query: str) -> bool:
        normalized = cls._normalize_query(query)
        if not cls._is_wfh_query(normalized):
            return False
        remote_specific = [
            "remote",
            "remotely",
            "permanent",
            "temporarily",
            "temporary",
            "consecutive",
            "week",
            "weeks",
            "month",
            "months",
        ]
        return not any(token in normalized for token in remote_specific)

    @classmethod
    def _doc_signal_features(cls, doc: Document) -> Dict[str, int]:
        metadata = doc.metadata or {}
        section_text = cls._normalize_text(metadata.get("section_path", ""))
        heading_text = cls._normalize_text(metadata.get("parent_heading", ""))
        content_text = cls._normalize_text(doc.page_content)
        full_text = " ".join(part for part in [section_text, heading_text, content_text] if part)

        wfh_signals = cls._configured_signal_hints("wfh")
        termination_signals = cls._configured_signal_hints("termination")
        resignation_signals = cls._configured_signal_hints("resignation")

        return {
            "wfh_primary": cls._count_markers(full_text, wfh_signals.get("primary_location", [])),
            "wfh_secondary": cls._count_markers(full_text, wfh_signals.get("secondary_location", [])),
            "wfh_allowance": cls._count_markers(full_text, wfh_signals.get("allowance", []))
            + cls._wfh_duration_hits(full_text),
            "wfh_approval": cls._count_markers(full_text, wfh_signals.get("approval", [])),
            "wfh_notice": cls._count_markers(full_text, wfh_signals.get("notice", [])),
            "wfh_conditions": cls._count_markers(full_text, wfh_signals.get("conditions", [])),
            "wfh_adjacent_noise": cls._count_markers(full_text, wfh_signals.get("adjacent_noise", [])),
            "termination_primary": cls._count_markers(full_text, termination_signals.get("primary", [])),
            "termination_discipline": cls._count_markers(full_text, termination_signals.get("discipline", [])),
            "termination_cause": cls._count_markers(full_text, termination_signals.get("cause", [])),
            "termination_resignation": cls._count_markers(full_text, termination_signals.get("resignation", [])),
            "termination_admin": cls._count_markers(full_text, termination_signals.get("admin", [])),
            "resignation_primary": cls._count_markers(full_text, resignation_signals.get("primary", [])),
            "resignation_secondary": cls._count_markers(full_text, resignation_signals.get("secondary", [])),
            "resignation_noise": cls._count_markers(full_text, resignation_signals.get("termination_noise", [])),
        }

    def _final_k_for_query(self, query: str) -> int:
        topic = self._topic_for_query(query)
        # Broad topics should include more context to avoid missing key subsections.
        if topic in {"wfh", "termination", "resignation"}:
            return min(int(self._final_k_broad), 4)
        if topic == "leave_timeoff":
            return int(self._final_k_broad)
        if topic != "focused":
            return int(self._final_k_broad)
        terms = self._query_terms(query)
        if len(terms) <= 2:
            return int(self._final_k_broad)
        return int(self._final_k_focused)

    @classmethod
    def _topic_for_query(cls, query: str) -> str:
        scores = cls._query_topic_scores(query)
        if not scores:
            q = (query or "").lower()
            if any(x in q for x in ["overtime", "working hours", "work hours", "payroll"]):
                return "hours_overtime"
            if any(x in q for x in ["benefit", "benefits", "perk", "perks", "insurance", "health"]):
                return "benefits"
            return "focused"

        top_topic = scores[0][0]
        reverse = {value: key for key, value in cls.TOPIC_ALIAS_MAP.items()}
        return reverse.get(top_topic, top_topic)

    @classmethod
    def _intent_hint_terms(cls, topic: str) -> List[str]:
        raw = os.environ.get("INTENT_HINTS_JSON")
        if raw:
            try:
                parsed = json.loads(raw)
                if isinstance(parsed, dict):
                    terms = parsed.get(topic) or parsed.get(topic.lower()) or []
                    if isinstance(terms, str):
                        terms = [terms]
                    if isinstance(terms, list):
                        cleaned = [cls._normalize_text(str(term)) for term in terms]
                        return [term for term in cleaned if term]
            except Exception:
                pass
        return list(cls.DEFAULT_INTENT_HINTS.get(topic, []))

    @classmethod
    def _create_qa_chain(cls):
        llm = AzureChatOpenAI(
            deployment_name=os.environ.get("AZURE_DEPLOYMENT_NAME", "gpt-35-turbo"),
            model_name=os.environ.get("AZURE_DEPLOYMENT_NAME", "gpt-35-turbo"),
            temperature=0,
        )
        return load_qa_chain(
            llm,
            chain_type="stuff",
            prompt=cls.QA_PROMPT,
            verbose=False,
        )

    @classmethod
    def load_docs(cls, directory: str):
        documents = []
        for filename in sorted(os.listdir(directory)):
            doc_path = os.path.join(directory, filename)
            if not os.path.isfile(doc_path):
                continue
            lower = filename.lower()
            ext = os.path.splitext(lower)[1]
            try:
                if lower.endswith(".docx"):
                    documents.extend(cls._load_docx(doc_path))
                elif lower.endswith(".pdf"):
                    documents.extend(cls._load_pdf(doc_path))
                elif ext in cls.GENERIC_TEXT_EXTENSIONS:
                    documents.extend(cls._load_text(doc_path))
                elif ext in (".pptx", ".xlsx"):
                    documents.extend(cls._load_office(doc_path))
                else:
                    print(f"Skipping unsupported file type: {doc_path}")
            except KeyError as e:
                print(f"Error loading {doc_path}: {e}")
            except Exception as e:
                print(f"Unexpected error loading {doc_path}: {e}")
        return documents

    @staticmethod
    def _normalize_text(text: str) -> str:
        return re.sub(r"\s+", " ", text or "").strip()

    @classmethod
    def _query_rewrites(cls) -> Dict[str, str]:
        rewrites = dict(cls.DEFAULT_QUERY_REWRITES)
        raw = os.environ.get("QUERY_NORMALIZATION_JSON")
        if raw:
            try:
                parsed = json.loads(raw)
                if isinstance(parsed, dict):
                    for src, dst in parsed.items():
                        src_text = cls._normalize_text(str(src)).lower()
                        dst_text = cls._normalize_text(str(dst)).lower()
                        if src_text and dst_text:
                            rewrites[src_text] = dst_text
            except Exception:
                pass
        return rewrites

    @classmethod
    def _normalize_query(cls, query: str) -> str:
        normalized = cls._normalize_text(query).lower()
        for src, dst in sorted(
            cls._query_rewrites().items(), key=lambda item: len(item[0]), reverse=True
        ):
            normalized = re.sub(rf"\b{re.escape(src)}\b", dst, normalized)
        return cls._normalize_text(normalized)

    @staticmethod
    def _safe_preview(text: str, limit: int = 120) -> str:
        return HelpCenterAgent._normalize_text(text)[:limit]

    @classmethod
    def _source_content(cls, doc: Document) -> str:
        metadata = doc.metadata or {}
        raw_content = metadata.get("raw_content")
        if isinstance(raw_content, str) and cls._normalize_text(raw_content):
            return raw_content
        return doc.page_content

    @classmethod
    def _query_terms(cls, query: str) -> List[str]:
        stop_words = {
            "the",
            "a",
            "an",
            "about",
            "tell",
            "me",
            "what",
            "is",
            "are",
            "policy",
            "policies",
            "regarding",
        }
        terms = re.findall(r"[a-zA-Z]+", cls._normalize_query(query))
        return [term for term in terms if term not in stop_words and len(term) > 2]

    @classmethod
    def _is_timeoff_leave_query(cls, query: str) -> bool:
        q = cls._normalize_query(query)
        # If user is clearly asking about separation/termination, don't route to time-off leave.
        separation_markers = [
            "termination",
            "terminate",
            "resign",
            "resignation",
            "quit",
            "notice",
            "layoff",
            "fired",
            "dismiss",
            "cobra",
        ]
        if any(marker in q for marker in separation_markers):
            return False
        timeoff_markers = [
            "leave",
            "holiday",
            "holidays",
            "pto",
            "sick",
            "bereavement",
            "jury duty",
            "voting",
            "parental",
            "maternity",
            "paternity",
            "illness",
            "time off",
        ]
        return any(marker in q for marker in timeoff_markers)

    @classmethod
    def _build_query_variants(cls, query: str, hints: Sequence[str]) -> List[str]:
        variants: List[str] = []
        for value in [query, cls._normalize_query(query), *hints]:
            normalized = cls._normalize_text(value)
            if normalized and normalized not in variants:
                variants.append(normalized)
            if len(variants) >= cls.MAX_QUERY_VARIANTS:
                break
        return variants

    @classmethod
    def _is_count_question(cls, query: str) -> bool:
        q = cls._normalize_query(query)
        return any(
            marker in q for marker in ["how many", "number of", "count ", "total "]
        )

    @classmethod
    def _is_holiday_count_query(cls, query: str) -> bool:
        q = cls._normalize_query(query)
        return cls._is_count_question(q) and any(
            marker in q for marker in ["holiday", "holidays"]
        )

    @classmethod
    def _extract_bracket_text(cls, text: str) -> str:
        match = re.search(r"\[([^\]]+)\]", text or "")
        if match:
            return cls._normalize_text(match.group(1))
        return cls._normalize_text(text)

    @classmethod
    def _load_source_blocks(cls, source_path: str) -> List[Document]:
        if source_path.lower().endswith(".pdf"):
            return cls._load_pdf(source_path)
        if source_path.lower().endswith(".docx"):
            return cls._load_docx(source_path)
        return []

    @classmethod
    def _holiday_names_from_sources(cls, matching_docs: Sequence[Document]) -> List[str]:
        holiday_names: List[str] = []
        seen = set()
        for source_path in {
            str((doc.metadata or {}).get("source", "")) for doc in matching_docs if doc.metadata
        }:
            if not source_path:
                continue
            try:
                blocks = cls._load_source_blocks(source_path)
            except Exception:
                continue
            for block in blocks:
                metadata = block.metadata or {}
                section_text = cls._normalize_text(
                    f"{metadata.get('section_path', '')} {metadata.get('parent_heading', '')}"
                ).lower()
                if "holiday" not in section_text:
                    continue
                raw_candidates: List[str] = []
                if metadata.get("block_type") == "bullet":
                    raw_candidates = [cls._extract_bracket_text(block.page_content)]
                elif (
                    metadata.get("block_type") == "paragraph"
                    and "observes the following holidays" in section_text
                ):
                    raw_candidates = [
                        cls._normalize_text(match)
                        for match in re.findall(r"\[([^\]]+)\]", block.page_content or "")
                    ]
                for candidate in raw_candidates:
                    candidate = re.sub(r"^([-\*\?]|\d+[\.\)])\s*", "", candidate).strip()
                    if not candidate:
                        continue
                    normalized = candidate.lower()
                    if normalized in seen:
                        continue
                    seen.add(normalized)
                    holiday_names.append(candidate)
        return holiday_names

    @classmethod
    def _best_sentence_for_signal(
        cls, sentences: Sequence[str], include_markers: Sequence[str], exclude_markers: Sequence[str] | None = None
    ) -> str:
        exclude_markers = exclude_markers or []
        best_sentence = ""
        best_score = -1
        for sentence in sentences:
            normalized = cls._normalize_text(sentence)
            if not normalized:
                continue
            lowered = normalized.lower()
            if exclude_markers and any(marker in lowered for marker in exclude_markers):
                continue
            score = cls._count_markers(lowered, include_markers) + cls._wfh_duration_hits(lowered)
            if score > best_score:
                best_sentence = normalized
                best_score = score
        return best_sentence if best_score > 0 else ""

    @classmethod
    def _wfh_answer_from_context(cls, matching_docs: Sequence[Document]) -> str | None:
        ranked_docs = sorted(
            matching_docs,
            key=lambda doc: (
                cls._doc_signal_features(doc)["wfh_primary"],
                cls._doc_signal_features(doc)["wfh_allowance"]
                + cls._doc_signal_features(doc)["wfh_approval"]
                + cls._doc_signal_features(doc)["wfh_notice"],
            ),
            reverse=True,
        )
        if not ranked_docs:
            return None

        primary_doc = ranked_docs[0]
        features = cls._doc_signal_features(primary_doc)
        if features["wfh_primary"] == 0 and features["wfh_secondary"] == 0:
            return None

        sentences = cls._split_sentences(cls._source_content(primary_doc))
        if not sentences:
            sentences = [cls._normalize_text(cls._source_content(primary_doc))]

        signal_hints = cls._configured_signal_hints("wfh")
        lead_markers = (
            signal_hints.get("primary_location", [])
            + signal_hints.get("allowance", [])
        )
        lead = cls._best_sentence_for_signal(sentences, lead_markers)
        approval = cls._best_sentence_for_signal(
            sentences,
            signal_hints.get("approval", []),
            exclude_markers=[lead.lower()] if lead else [],
        )
        notice = cls._best_sentence_for_signal(
            sentences,
            signal_hints.get("notice", []),
            exclude_markers=[value.lower() for value in [lead, approval] if value],
        )
        conditions = cls._best_sentence_for_signal(
            sentences,
            signal_hints.get("conditions", []),
            exclude_markers=[value.lower() for value in [lead, approval, notice] if value],
        )

        bullets: List[str] = []
        for sentence in [lead, approval, notice, conditions]:
            cleaned = cls._normalize_text(sentence)
            if cleaned and cleaned not in bullets:
                bullets.append(cleaned.rstrip(".") + ".")

        if not bullets:
            return None
        return "\n".join(f"- {bullet}" for bullet in bullets)

    @classmethod
    def _policy_primary_markers(cls) -> List[str]:
        return [
            "allow",
            "allowed",
            "can",
            "may",
            "normally",
            "eligible",
            "maximum",
            "minimum",
            "up to",
            "per week",
            "per month",
            "if your job",
        ]

    @classmethod
    def _policy_operational_markers(cls) -> List[str]:
        return [
            "manager",
            "hr",
            "request",
            "inform",
            "approval",
            "approve",
            "advance",
            "secure",
            "security",
            "email",
            "call",
            "check in",
        ]

    @classmethod
    def _should_use_ordered_policy_summary(
        cls, normalized_query: str, matching_docs: Sequence[Document]
    ) -> bool:
        if not matching_docs:
            return False
        top_doc = matching_docs[0]
        primary_topic = str((top_doc.metadata or {}).get("policy_focus", ""))
        if primary_topic != "remote_work":
            return False
        if not cls._is_wfh_query(normalized_query):
            return False
        broad_markers = [
            "details",
            "summary",
            "how much",
            "bonus",
            "regarding",
            "all about",
        ]
        if any(marker in normalized_query for marker in broad_markers):
            return False
        return True

    @classmethod
    def _ordered_policy_summary_from_context(
        cls, normalized_query: str, matching_docs: Sequence[Document]
    ) -> str | None:
        if not matching_docs:
            return None
        if not cls._should_use_ordered_policy_summary(normalized_query, matching_docs):
            return None

        top_doc = matching_docs[0]
        primary_topic = str((top_doc.metadata or {}).get("policy_focus", ""))
        if not primary_topic:
            return None

        related_docs = [
            doc for doc in matching_docs[:3] if str((doc.metadata or {}).get("policy_focus", "")) == primary_topic
        ]
        if not related_docs:
            related_docs = [top_doc]

        combined_text = " ".join(cls._source_content(doc) for doc in related_docs)
        sentences = cls._split_sentences(combined_text)
        if len(sentences) < 2:
            return None

        primary_sentence = cls._best_sentence_for_signal(
            sentences, cls._policy_primary_markers()
        )
        operational_sentence = cls._best_sentence_for_signal(
            sentences,
            cls._policy_operational_markers(),
            exclude_markers=[primary_sentence.lower()] if primary_sentence else [],
        )

        if not primary_sentence or not operational_sentence:
            return None

        primary_hits = cls._count_markers(primary_sentence, cls._query_terms(normalized_query))
        topic_hits = cls._count_markers(primary_sentence, cls.DEFAULT_POLICY_TOPIC_HINTS.get(primary_topic, []))
        if primary_hits == 0 and topic_hits == 0:
            return None

        ordered: List[str] = []
        for sentence in [primary_sentence, operational_sentence]:
            cleaned = cls._normalize_text(sentence)
            if cleaned and cleaned not in ordered:
                ordered.append(cleaned.rstrip(".") + ".")

        extra_sentences = []
        for sentence in sentences:
            cleaned = cls._normalize_text(sentence)
            if (
                cleaned
                and cleaned not in [primary_sentence, operational_sentence]
                and (
                    cls._count_markers(cleaned, cls._policy_operational_markers()) > 0
                    or cls._count_markers(cleaned, cls._query_terms(normalized_query)) > 0
                )
            ):
                extra_sentences.append(cleaned.rstrip(".") + ".")
            if len(extra_sentences) >= 2:
                break

        for sentence in extra_sentences:
            if sentence not in ordered:
                ordered.append(sentence)

        return "\n".join(f"- {sentence}" for sentence in ordered) if ordered else None

    @classmethod
    def _answer_from_context(
        cls, original_query: str, normalized_query: str, matching_docs: Sequence[Document]
    ) -> str | None:
        if cls._is_holiday_count_query(normalized_query):
            holiday_names = cls._holiday_names_from_sources(matching_docs)
            if holiday_names:
                joined = ", ".join(holiday_names)
                return (
                    f"The handbook lists {len(holiday_names)} company holidays: {joined}."
                )
        ordered_policy_summary = cls._ordered_policy_summary_from_context(
            normalized_query, matching_docs
        )
        if ordered_policy_summary is not None:
            return ordered_policy_summary
        return None

    @classmethod
    def _count_term_hits(cls, text: str, terms: Sequence[str]) -> int:
        return sum(1 for term in terms if cls._marker_present(text, term))

    @classmethod
    def _candidate_rank(cls, query: str, doc: Document, score: float) -> float:
        metadata = doc.metadata or {}
        terms = cls._query_terms(query)
        topic = cls._topic_for_query(query)
        normalized_query = cls._normalize_query(query)
        intent_terms = cls._intent_hint_terms(topic)
        if intent_terms:
            terms = sorted(set(terms + intent_terms))
        section_text = cls._normalize_text(metadata.get("section_path", "")).lower()
        heading_text = cls._normalize_text(metadata.get("parent_heading", "")).lower()
        heading_hits = cls._count_term_hits(heading_text, terms)
        section_hits = cls._count_term_hits(section_text, terms)
        content_hits = cls._count_term_hits(doc.page_content[:400], terms)
        body_bonus = 0.04 if metadata.get("block_type") == "paragraph" else 0.0
        features = cls._doc_signal_features(doc)

        penalty = 0.0
        bonus = 0.0
        if topic == "leave_timeoff":
            # These sections contain the word "leave"/"leaving" but are not time-off leave policy.
            if "leaving our company" in section_text or "leaving our company" in heading_text:
                penalty += 0.35
            if "cobra" in section_text or "cobra" in heading_text:
                penalty += 0.35
            if "termination" in section_text or "termination" in heading_text:
                penalty += 0.18
            if "progressive discipline" in section_text or "progressive discipline" in heading_text:
                penalty += 0.12
            if any(term in normalized_query for term in ["sick", "illness", "fmla"]):
                if "sick leave" in section_text or "sick leave" in heading_text:
                    bonus += 0.22
                if "long-term illness" in section_text or "long-term illness" in heading_text:
                    bonus += 0.26
            if "holiday" in normalized_query:
                if "holiday" in section_text or "holiday" in heading_text:
                    bonus += 0.24
            if "pto" in normalized_query or "paid time off" in normalized_query:
                if "paid time off" in section_text or "paid time off" in heading_text:
                    bonus += 0.18
            if "leave" in normalized_query and (
                "bereavement leave" in section_text
                or "parental leave" in section_text
                or "jury duty" in section_text
            ):
                bonus += 0.08
            if metadata.get("block_type") == "bullet" and "holiday" in section_text:
                bonus += 0.1
            if "time" == section_text:
                penalty += 0.08
            if "prescription drugs" in section_text:
                penalty += 0.1
        elif topic == "wfh":
            bonus += 0.26 * features["wfh_primary"]
            bonus += 0.14 * features["wfh_secondary"]
            bonus += 0.12 * features["wfh_allowance"]
            bonus += 0.10 * features["wfh_approval"]
            bonus += 0.10 * features["wfh_notice"]
            bonus += 0.05 * features["wfh_conditions"]
            penalty += 0.10 * features["wfh_adjacent_noise"]
            if cls._is_simple_wfh_rules_query(normalized_query):
                if features["wfh_secondary"] > 0 and features["wfh_primary"] == 0:
                    penalty += 0.10
                if features["wfh_primary"] == 0 and features["wfh_adjacent_noise"] > 0:
                    penalty += 0.18
        elif topic == "termination":
            bonus += 0.26 * features["termination_primary"]
            bonus += 0.12 * features["termination_discipline"]
            bonus += 0.10 * features["termination_cause"]
            penalty += 0.10 * features["termination_admin"]
            if not cls._is_resignation_query(normalized_query):
                penalty += 0.16 * features["termination_resignation"]
                if (
                    features["termination_resignation"] > 0
                    and features["termination_cause"] + features["termination_discipline"] == 0
                ):
                    penalty += 0.12
        elif topic == "resignation":
            bonus += 0.26 * features["resignation_primary"]
            bonus += 0.08 * features["resignation_secondary"]
            penalty += 0.14 * features["resignation_noise"]
        elif topic not in {"termination", "resignation"}:
            if "cobra" in section_text or "cobra" in heading_text:
                penalty += 0.18
            if "termination" in section_text or "termination" in heading_text:
                penalty += 0.12

        return (
            float(score)
            - (0.14 * heading_hits)
            - (0.08 * section_hits)
            - (0.04 * content_hits)
            - body_bonus
            - bonus
            + penalty
        )

    @staticmethod
    def _is_heading(text: str, style_name: str = "") -> bool:
        normalized = HelpCenterAgent._normalize_text(text)
        if not normalized:
            return False
        if style_name.lower().startswith("heading"):
            return True
        if len(normalized) <= 120:
            if re.match(r"^\d+(\.\d+)*\s+[A-Za-z].*", normalized):
                return True
            if re.match(r"^[A-Z][A-Za-z/&,\-\s\(\)]+?\s+\d+$", normalized):
                return True
            if normalized.isupper() and len(normalized.split()) <= 12:
                return True
            if (
                normalized[0].isupper()
                and not normalized.endswith(".")
                and len(normalized.split()) <= 10
            ):
                return True
        return False

    @staticmethod
    def _heading_level(style_name: str) -> int:
        match = re.search(r"heading\s*(\d+)", style_name.lower())
        if match:
            return max(1, min(6, int(match.group(1))))
        if style_name.lower().startswith("heading"):
            return 2
        return 1

    @staticmethod
    def _is_bullet_line(text: str) -> bool:
        return bool(
            re.match(r"^([-\*]|\d+[\.\)])\s+", HelpCenterAgent._normalize_text(text))
        )

    @staticmethod
    def _is_table_like(text: str) -> bool:
        normalized = HelpCenterAgent._normalize_text(text)
        if "|" in normalized:
            return True
        return bool(re.search(r"\s{2,}", normalized) and len(normalized.split()) <= 18)

    @staticmethod
    def _looks_like_toc_line(text: str) -> bool:
        normalized = HelpCenterAgent._normalize_text(text)
        return bool(
            normalized
            and len(normalized.split()) <= 10
            and re.match(r"^[A-Za-z][A-Za-z/&,\-\s\(\)]+?\s+\d+$", normalized)
        )

    @classmethod
    def _is_probable_pdf_toc_page(cls, lines: Sequence[str], page_number: int) -> bool:
        # NOTE: This only affects indexing. If you change TOC_PAGE_SCAN_LIMIT,
        # rebuild the Chroma index (or bump CHUNKING_VERSION).
        try:
            limit = int(os.environ.get("TOC_PAGE_SCAN_LIMIT", cls.TOC_PAGE_SCAN_LIMIT))
        except ValueError:
            limit = int(cls.TOC_PAGE_SCAN_LIMIT)
        if page_number + 1 > limit:
            return False
        normalized = [cls._normalize_text(line) for line in lines if cls._normalize_text(line)]
        if len(normalized) < 8:
            return False
        toc_lines = sum(1 for line in normalized if cls._looks_like_toc_line(line))
        long_lines = sum(1 for line in normalized if len(line.split()) >= 12)
        return toc_lines >= 8 and toc_lines >= max(1, len(normalized) // 2) and long_lines <= 3

    @classmethod
    def _classify_block(cls, text: str, style_name: str = "") -> str:
        if cls._is_heading(text, style_name):
            return "heading"
        if cls._is_bullet_line(text):
            return "bullet"
        if cls._is_table_like(text):
            return "table_like"
        return "paragraph"

    @staticmethod
    def _iter_docx_items(doc: DocxNativeDocument):
        for child in doc.element.body.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, doc)
            elif isinstance(child, CT_Tbl):
                yield Table(child, doc)

    @classmethod
    def _docx_section_path(cls, section_stack: Sequence[Tuple[int, str]]) -> str:
        return " > ".join(title for _, title in section_stack)

    @classmethod
    def _build_metadata(
        cls,
        source: str,
        document_type: str,
        section_stack: Sequence[Tuple[int, str]],
        extra: Dict[str, object] | None = None,
    ) -> Dict[str, object]:
        metadata: Dict[str, object] = {
            "source": source,
            "document_type": document_type,
            "section_path": cls._docx_section_path(section_stack),
            "section_depth": len(section_stack),
        }
        if extra:
            metadata.update(extra)
        return metadata

    # Extensions handled by the generic text loader (everything that is not
    # .docx/.pdf but can be read as text). The UI upload allowlist is the
    # authoritative gate; this is the indexing-side safety net.
    GENERIC_TEXT_EXTENSIONS = (
        ".txt",
        ".md",
        ".markdown",
        ".rst",
        ".csv",
        ".tsv",
        ".json",
        ".html",
        ".htm",
        ".log",
        ".text",
        ".rtf",
    )
    MAX_TEXT_FILE_BYTES = 10 * 1024 * 1024  # 10 MB guard for plain-text loads

    @staticmethod
    def _strip_html(raw: str) -> str:
        # Drop scripts/styles, then tags, then unescape a few common entities.
        without_blocks = re.sub(
            r"<(script|style)[^>]*>.*?</\1>", " ", raw, flags=re.IGNORECASE | re.DOTALL
        )
        # Turn block-level closers into newlines so paragraphs survive.
        with_breaks = re.sub(
            r"</(p|div|li|tr|h[1-6]|br|section|article)\s*>",
            "\n",
            without_blocks,
            flags=re.IGNORECASE,
        )
        with_breaks = re.sub(r"<br\s*/?>", "\n", with_breaks, flags=re.IGNORECASE)
        no_tags = re.sub(r"<[^>]+>", " ", with_breaks)
        replacements = {
            "&nbsp;": " ",
            "&amp;": "&",
            "&lt;": "<",
            "&gt;": ">",
            "&quot;": '"',
            "&#39;": "'",
            "&apos;": "'",
        }
        for entity, char in replacements.items():
            no_tags = no_tags.replace(entity, char)
        return no_tags

    @classmethod
    def _text_to_blocks(cls, file_path: str, text: str, source_type: str, merge_short_lines: bool = False) -> List[Document]:
        blocks: List[Document] = []
        section_stack: List[Tuple[int, str]] = []
        short_line_buffer: List[str] = []
        # Track the very first heading seen — used as doc-level context so
        # merged bullet chunks carry the document topic ("Airport Lounge Access
        # Policy Facilities: …") and score well for topic queries.
        doc_title: List[str] = []  # list so closure can mutate it

        def flush_short_buffer() -> None:
            if not short_line_buffer:
                return
            # Prefix = "<doc title> <current section>:" gives natural,
            # reranker-friendly text that includes the document topic.
            title_part = doc_title[0] if doc_title else ""
            section_part = (
                cls._normalize_text(section_stack[-1][1]) if section_stack else ""
            )
            if title_part and title_part != section_part:
                prefix = f"{title_part} {section_part}: "
            elif section_part:
                prefix = f"{section_part}: "
            else:
                prefix = ""
            merged = cls._normalize_text(prefix + " ".join(short_line_buffer))
            short_line_buffer.clear()
            if not merged:
                return
            blocks.append(
                Document(
                    page_content=merged,
                    metadata=cls._build_metadata(
                        file_path,
                        source_type,
                        section_stack,
                        {"block_type": "bullet"},
                    ),
                )
            )

        raw_paragraphs = re.split(r"\n\s*\n", text)
        for raw_paragraph in raw_paragraphs:
            for raw_line in raw_paragraph.split("\n"):
                stripped = raw_line.strip()
                if not stripped:
                    continue

                # Markdown ATX headings flush buffer + update section context.
                md_heading = re.match(r"^(#{1,6})\s+(.*)$", stripped)
                if md_heading:
                    flush_short_buffer()
                    level = len(md_heading.group(1))
                    title = cls._normalize_text(md_heading.group(2))
                    if title:
                        if not doc_title:
                            doc_title.append(title)
                        section_stack = [e for e in section_stack if e[0] < level]
                        section_stack.append((level, title))
                        blocks.append(
                            Document(
                                page_content=title,
                                metadata=cls._build_metadata(
                                    file_path,
                                    source_type,
                                    section_stack,
                                    {
                                        "block_type": "heading",
                                        "heading_level": level,
                                        "is_heading_only": True,
                                    },
                                ),
                            )
                        )
                    continue

                normalized = cls._normalize_text(stripped)
                if not normalized:
                    continue

                # HTML bullet-merge: in this mode, short lines bypass heading
                # detection entirely — they're <li> items, not section headings.
                # MUST come before _is_heading to prevent false positives.
                if merge_short_lines and len(normalized) < cls.MIN_BODY_CHARS:
                    short_line_buffer.append(normalized)
                    continue

                # List-intro lines (end with ":") seed the buffer so that the
                # intro + following bullets become one combined chunk.
                if merge_short_lines and normalized.endswith(":"):
                    flush_short_buffer()
                    short_line_buffer.append(normalized)
                    continue

                # Normal heading detection — only reached for longer lines
                # (< MIN_BODY_CHARS already handled above in HTML mode).
                if cls._is_heading(normalized):
                    flush_short_buffer()
                    level = 2
                    if not doc_title:
                        doc_title.append(normalized)
                    section_stack = [e for e in section_stack if e[0] < level]
                    section_stack.append((level, normalized))
                    blocks.append(
                        Document(
                            page_content=normalized,
                            metadata=cls._build_metadata(
                                file_path,
                                source_type,
                                section_stack,
                                {
                                    "block_type": "heading",
                                    "heading_level": level,
                                    "is_heading_only": True,
                                },
                            ),
                        )
                    )
                    continue

                # Substantive long line — flush any buffer first, then emit.
                flush_short_buffer()
                blocks.append(
                    Document(
                        page_content=normalized,
                        metadata=cls._build_metadata(
                            file_path,
                            source_type,
                            section_stack,
                            {"block_type": cls._classify_block(normalized)},
                        ),
                    )
                )

        flush_short_buffer()  # emit any remaining buffered lines
        return blocks

    @classmethod
    def _flatten_json_to_blocks(
        cls, file_path: str, obj: object, _section_stack: List[Tuple[int, str]] | None = None
    ) -> List[Document]:
        """Recursively flatten a parsed JSON object into readable Document blocks.

        Each dict item becomes a key: value sentence; each list of dicts becomes
        one block per list item with all its fields joined. This avoids the JSON
        syntax noise ({, [, "key":) that confuses the short-line merge buffer.
        """
        if _section_stack is None:
            _section_stack = []
        blocks: List[Document] = []

        def make_doc(text: str, stack: List[Tuple[int, str]]) -> None:
            normalized = cls._normalize_text(text)
            if not normalized:
                return
            blocks.append(
                Document(
                    page_content=normalized,
                    metadata=cls._build_metadata(
                        file_path, "json", stack, {"block_type": "paragraph"}
                    ),
                )
            )

        def flatten_value(v: object) -> str:
            """Render a scalar or simple value as a string."""
            if isinstance(v, str):
                return cls._normalize_text(v)
            if isinstance(v, (int, float, bool)):
                return str(v)
            if isinstance(v, list):
                parts = [flatten_value(item) for item in v if not isinstance(item, (dict, list))]
                return ", ".join(p for p in parts if p)
            return ""

        def process(node: object, stack: List[Tuple[int, str]]) -> None:
            if isinstance(node, dict):
                # If a dict has a "name" or "title" key, use it as a heading.
                name_key = next(
                    (k for k in ("name", "title", "id", "type") if k in node), None
                )
                if name_key:
                    heading = cls._normalize_text(str(node[name_key]))
                    new_stack = [e for e in stack if e[0] < 2] + [(2, heading)]
                else:
                    new_stack = stack

                # Collect all scalar fields into one sentence.
                parts: List[str] = []
                for k, v in node.items():
                    if isinstance(v, (dict, list)):
                        continue
                    val = flatten_value(v)
                    if val:
                        parts.append(f"{cls._normalize_text(str(k))}: {val}")
                if parts:
                    make_doc(" | ".join(parts), new_stack)

                # Recurse into nested structures.
                for v in node.values():
                    if isinstance(v, (dict, list)):
                        process(v, new_stack)

            elif isinstance(node, list):
                for item in node:
                    process(item, stack)

            elif isinstance(node, str):
                val = cls._normalize_text(node)
                if val and len(val) >= cls.MIN_BODY_CHARS:
                    make_doc(val, stack)

        process(obj, _section_stack)
        return blocks

    @classmethod
    def _load_text(cls, file_path):
        """Generic loader for text-extractable files (txt/md/csv/json/html/...)."""
        try:
            try:
                if os.path.getsize(file_path) > cls.MAX_TEXT_FILE_BYTES:
                    print(f"Skipping oversized text file: {file_path}")
                    return []
            except OSError:
                pass

            with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
                raw = handle.read()

            ext = os.path.splitext(file_path)[1].lower()

            if ext in (".html", ".htm"):
                text = cls._strip_html(raw)
            elif ext in (".csv", ".tsv"):
                # CSV/TSV: build Documents directly so rows aren't split back
                # into individual short lines by _text_to_blocks.
                # block_type="table_row" bypasses MIN_BODY_CHARS/WORDS checks.
                delimiter = "\t" if ext == ".tsv" else ","
                rows = []
                for line in raw.splitlines():
                    if not line.strip():
                        continue
                    cells = [cell.strip() for cell in line.split(delimiter)]
                    row_text = " | ".join(cell for cell in cells if cell)
                    if row_text:
                        rows.append(row_text)
                header = rows[0] if rows else ""
                data_rows = rows[1:] if len(rows) > 1 else rows
                CSV_CHUNK = 5
                csv_section: List[Tuple[int, str]] = [(1, header)] if header else []
                csv_blocks: List[Document] = []
                for _i in range(0, max(len(data_rows), 1), CSV_CHUNK):
                    group = data_rows[_i: _i + CSV_CHUNK]
                    chunk_text = cls._normalize_text(
                        (header + "\n" if header else "") + "\n".join(group)
                    )
                    if not chunk_text:
                        continue
                    csv_blocks.append(
                        Document(
                            page_content=chunk_text,
                            metadata=cls._build_metadata(
                                file_path,
                                "csv",
                                csv_section,
                                {"block_type": "table_row"},
                            ),
                        )
                    )
                return csv_blocks
            elif ext == ".json":
                # Flatten JSON into readable sentences directly — bypass
                # _text_to_blocks so JSON syntax noise never enters the index.
                try:
                    parsed = json.loads(raw)
                    json_blocks = cls._flatten_json_to_blocks(
                        file_path, parsed
                    )
                    if json_blocks:
                        return json_blocks
                    text = json.dumps(parsed, indent=2, ensure_ascii=False)
                except Exception:
                    text = raw
            else:
                text = raw

            source_type = {
                ".csv": "csv",
                ".tsv": "csv",
                ".json": "json",
                ".html": "html",
                ".htm": "html",
                ".md": "markdown",
                ".markdown": "markdown",
            }.get(ext, "text")

            # Pass merge_short_lines=True only for HTML so <li> bullets
            # are merged into substantive chunks; TXT/MD don't need it.
            merge_bullets = ext in (".html", ".htm")
            return cls._text_to_blocks(file_path, text, source_type, merge_short_lines=merge_bullets)
        except Exception as e:
            raise Exception(f"Failed to load text file: {e}")

    @classmethod
    def _load_office(cls, file_path):
        """Optional loader for .pptx/.xlsx — only if the libs are installed."""
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".pptx":
                from pptx import Presentation  # type: ignore

                prs = Presentation(file_path)
                lines: List[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False):
                            for paragraph in shape.text_frame.paragraphs:
                                line = cls._normalize_text(
                                    "".join(run.text for run in paragraph.runs)
                                )
                                if line:
                                    lines.append(line)
                return cls._text_to_blocks(file_path, "\n".join(lines), "pptx")
            if ext == ".xlsx":
                from openpyxl import load_workbook  # type: ignore

                wb = load_workbook(file_path, read_only=True, data_only=True)
                rows: List[str] = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        cells = [str(c).strip() for c in row if c is not None]
                        row_text = " | ".join(cell for cell in cells if cell)
                        if row_text:
                            rows.append(row_text)
                return cls._text_to_blocks(file_path, "\n".join(rows), "xlsx")
        except ImportError:
            print(f"Skipping {file_path}: optional library for {ext} not installed")
            return []
        except Exception as e:
            print(f"Unexpected error loading {file_path}: {e}")
            return []
        return []

    @classmethod
    def _load_docx(cls, file_path):
        try:
            doc = DocxDocument(file_path)
            blocks = []
            section_stack: List[Tuple[int, str]] = []
            table_index = 0

            for item in cls._iter_docx_items(doc):
                if isinstance(item, Paragraph):
                    text = cls._normalize_text(item.text)
                    if not text:
                        continue
                    style_name = getattr(getattr(item, "style", None), "name", "") or ""

                    if cls._is_heading(text, style_name):
                        level = cls._heading_level(style_name)
                        section_stack = [
                            entry for entry in section_stack if entry[0] < level
                        ]
                        section_stack.append((level, text))
                        blocks.append(
                            Document(
                                page_content=text,
                                metadata=cls._build_metadata(
                                    file_path,
                                    "docx",
                                    section_stack,
                                    {
                                        "block_type": "heading",
                                        "heading_level": level,
                                        "is_heading_only": True,
                                    },
                                ),
                            )
                        )
                        continue

                    blocks.append(
                        Document(
                            page_content=text,
                            metadata=cls._build_metadata(
                                file_path,
                                "docx",
                                section_stack,
                                {
                                    "block_type": cls._classify_block(text, style_name),
                                    "paragraph_style": style_name,
                                },
                            ),
                        )
                    )

                elif isinstance(item, Table):
                    table_index += 1
                    for row_index, row in enumerate(item.rows, start=1):
                        row_content = [cls._normalize_text(cell.text) for cell in row.cells]
                        row_text = cls._normalize_text(
                            " | ".join(part for part in row_content if part)
                        )
                        if not row_text:
                            continue
                        blocks.append(
                            Document(
                                page_content=row_text,
                                metadata=cls._build_metadata(
                                    file_path,
                                    "docx",
                                    section_stack,
                                    {
                                        "block_type": "table_row",
                                        "table_index": table_index,
                                        "row_index": row_index,
                                    },
                                ),
                            )
                        )

            return blocks
        except KeyError as e:
            raise KeyError(
                f"There is no item named '[Content_Types].xml' in the archive: {e}"
            )
        except Exception as e:
            raise Exception(f"Failed to load .docx file: {e}")

    @classmethod
    def _load_pdf(cls, file_path):
        try:
            loader = PyPDFLoader(file_path)
            pages = loader.load()
            blocks = []
            section_stack: List[Tuple[int, str]] = []
            for page_doc in pages:
                page_number = page_doc.metadata.get("page", 0)
                page_blocks, section_stack = cls._split_pdf_page_into_blocks(
                    file_path, page_doc, page_number, section_stack
                )
                blocks.extend(page_blocks)
            return blocks
        except Exception as e:
            raise Exception(f"Failed to load .pdf file: {e}")

    @classmethod
    def _split_pdf_page_into_blocks(
        cls,
        file_path: str,
        page_doc: Document,
        page_number: int,
        section_stack: Sequence[Tuple[int, str]],
    ):
        raw_lines = page_doc.page_content.splitlines()
        if cls._is_probable_pdf_toc_page(raw_lines, page_number):
            return [], list(section_stack)

        blocks = []
        current_sections: List[Tuple[int, str]] = list(section_stack)
        paragraph_buffer: List[str] = []
        paragraph_index = 0

        def flush_paragraph_buffer():
            nonlocal paragraph_index
            if not paragraph_buffer:
                return
            paragraph_index += 1
            paragraph_text = cls._normalize_text(" ".join(paragraph_buffer))
            paragraph_buffer.clear()
            if not paragraph_text:
                return
            blocks.append(
                Document(
                    page_content=paragraph_text,
                    metadata=cls._build_metadata(
                        file_path,
                        "pdf",
                        current_sections,
                        {
                            "block_type": "paragraph",
                            "page_number": page_number + 1,
                            "paragraph_index": paragraph_index,
                        },
                    ),
                )
            )

        for raw_line in raw_lines:
            line = cls._normalize_text(raw_line)
            if not line:
                flush_paragraph_buffer()
                continue

            if cls._is_heading(line):
                flush_paragraph_buffer()
                level = cls._heading_level("")
                current_sections = [entry for entry in current_sections if entry[0] < level]
                current_sections.append((level, line))
                blocks.append(
                    Document(
                        page_content=line,
                        metadata=cls._build_metadata(
                            file_path,
                            "pdf",
                            current_sections,
                            {
                                "block_type": "heading",
                                "page_number": page_number + 1,
                                "heading_level": level,
                                "is_heading_only": True,
                            },
                        ),
                    )
                )
                continue

            if cls._is_bullet_line(line):
                flush_paragraph_buffer()
                blocks.append(
                    Document(
                        page_content=line,
                        metadata=cls._build_metadata(
                            file_path,
                            "pdf",
                            current_sections,
                            {"block_type": "bullet", "page_number": page_number + 1},
                        ),
                    )
                )
                continue

            if cls._is_table_like(line):
                flush_paragraph_buffer()
                blocks.append(
                    Document(
                        page_content=line,
                        metadata=cls._build_metadata(
                            file_path,
                            "pdf",
                            current_sections,
                            {"block_type": "table_like", "page_number": page_number + 1},
                        ),
                    )
                )
                continue

            paragraph_buffer.append(line)

        flush_paragraph_buffer()
        return blocks, current_sections

    def _get_semantic_model(self) -> SentenceTransformer:
        if self._semantic_model is None:
            self._semantic_model = SentenceTransformer(self.EMBEDDING_MODEL_NAME)
        return self._semantic_model

    @staticmethod
    def _split_sentences(text: str) -> List[str]:
        return [part.strip() for part in re.split(r"(?<=[.!?])\s+", text.strip()) if part.strip()]

    @classmethod
    def _copy_metadata(cls, metadata: Dict[str, object], **updates) -> Dict[str, object]:
        copied = dict(metadata or {})
        copied.update(updates)
        return copied

    @classmethod
    def _attach_heading_context(
        cls, metadata: Dict[str, object], heading_stack: Sequence[Tuple[int, str]]
    ) -> Dict[str, object]:
        heading_context = " > ".join(title for _, title in heading_stack)
        parent_heading = heading_stack[-1][1] if heading_stack else ""
        return cls._copy_metadata(
            metadata,
            heading_context=heading_context,
            parent_heading=parent_heading,
            is_heading_only=False,
        )

    @classmethod
    def _chunk_text_with_recursive_fallback(cls, text: str) -> List[str]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=cls.LONG_SENTENCE_SPLIT_CHUNK_SIZE,
            chunk_overlap=cls.LONG_SENTENCE_SPLIT_CHUNK_OVERLAP,
        )
        return splitter.split_text(text)

    @classmethod
    def _has_substantive_body(cls, text: str, metadata: Dict[str, object]) -> bool:
        if metadata.get("block_type") == "table_row":
            return True
        normalized = cls._normalize_text(text)
        return len(normalized) >= cls.MIN_BODY_CHARS and len(normalized.split()) >= cls.MIN_BODY_WORDS

    @classmethod
    def _semantic_chunk_text(cls, raw_text: str, metadata: Dict[str, object]) -> str:
        raw_text = cls._normalize_text(raw_text)
        heading_context = cls._normalize_text(metadata.get("heading_context", ""))
        parent_heading = cls._normalize_text(metadata.get("parent_heading", ""))
        section_path = cls._normalize_text(metadata.get("section_path", ""))
        context_parts: List[str] = []

        if parent_heading:
            context_parts.append(f"Section: {parent_heading}.")
        elif section_path:
            context_parts.append(f"Section: {section_path}.")

        if heading_context and heading_context not in {parent_heading, section_path}:
            context_parts.append(f"Context: {heading_context}.")

        if raw_text:
            context_parts.append(raw_text)
        return " ".join(context_parts).strip()

    @classmethod
    def _is_atomic_block(cls, block_type: str) -> bool:
        return block_type in {"bullet", "table_like", "table_row"}

    @classmethod
    def _policy_topic_flags(cls, text: str, metadata: Dict[str, object]) -> Dict[str, str]:
        structure_text = cls._normalize_text(
            " ".join(
                [
                    cls._normalize_text(metadata.get("section_path", "")),
                    cls._normalize_text(metadata.get("parent_heading", "")),
                    cls._normalize_text(metadata.get("heading_context", "")),
                ]
            )
        ).lower()
        content_text = cls._normalize_text(text).lower()
        flags: Dict[str, str] = {}
        best_topic = ""
        best_score = 0
        for topic, markers in cls.DEFAULT_POLICY_TOPIC_HINTS.items():
            policy_markers = list(markers)
            legacy = cls._legacy_topic_name(topic)
            policy_markers.extend(cls._intent_hint_terms(legacy))
            structure_hits = cls._count_markers(structure_text, policy_markers)
            content_hits = cls._count_markers(content_text, policy_markers)
            preferred_hits = cls._count_markers(
                f"{structure_text} {content_text}",
                cls.TOPIC_PREFERRED_SECTION_MARKERS.get(topic, []),
            )
            avoid_hits = cls._count_markers(
                f"{structure_text} {content_text}",
                cls.TOPIC_AVOID_SECTION_MARKERS.get(topic, []),
            )
            score = (2 * structure_hits) + content_hits + preferred_hits - (2 * avoid_hits)
            is_topic = (
                (structure_hits > 0 and avoid_hits == 0)
                or (
                    score >= cls.MIN_TOPIC_FLAG_SCORE
                    and preferred_hits > avoid_hits
                )
                or (score >= cls.MIN_TOPIC_FLAG_SCORE + 1 and avoid_hits == 0)
            )
            flags[f"topic_{topic}"] = "1" if is_topic else "0"
            if is_topic and score > best_score:
                best_score = score
                best_topic = topic
        if best_topic and best_score >= cls.MIN_TOPIC_FLAG_SCORE:
            flags["policy_focus"] = best_topic
        return flags

    @classmethod
    def _child_docs_from_parent(cls, doc: Document) -> List[Document]:
        metadata = dict(doc.metadata or {})
        raw_content = cls._source_content(doc)
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=cls.CHILD_CHUNK_SIZE,
            chunk_overlap=cls.CHILD_CHUNK_OVERLAP,
        )
        pieces = [
            cls._normalize_text(piece)
            for piece in splitter.split_text(raw_content)
            if cls._normalize_text(piece)
        ]
        if not pieces:
            pieces = [cls._normalize_text(raw_content)]

        children: List[Document] = []
        for index, piece in enumerate(pieces, start=1):
            child_metadata = cls._copy_metadata(
                metadata,
                raw_content=piece,
                child_index=index,
                child_count=len(pieces),
                chunk_stage="child_chunk",
                chunk_strategy="parent_child_semantic",
            )
            child_text = cls._semantic_chunk_text(piece, child_metadata)
            children.append(Document(page_content=child_text, metadata=child_metadata))
        return children

    @classmethod
    def _semantic_seed_segments(cls, docs: Sequence[Document]) -> List[List[Document]]:
        windows: List[List[Document]] = []
        current_window: List[Document] = []
        current_chars = 0
        current_sentences = 0

        def flush_window() -> None:
            nonlocal current_window, current_chars, current_sentences
            if current_window:
                windows.append(current_window)
            current_window = []
            current_chars = 0
            current_sentences = 0

        for doc in docs:
            raw_text = cls._normalize_text(doc.page_content)
            if not raw_text or not cls._has_substantive_body(raw_text, doc.metadata or {}):
                continue

            block_type = str((doc.metadata or {}).get("block_type", "paragraph"))
            sentence_count = max(1, len(cls._split_sentences(raw_text)))
            should_be_atomic = cls._is_atomic_block(block_type)

            if should_be_atomic:
                flush_window()
                atomic_parts = [raw_text]
                if len(raw_text) > cls.MAX_SEMANTIC_WINDOW_CHARS:
                    atomic_parts = [
                        cls._normalize_text(part)
                        for part in cls._chunk_text_with_recursive_fallback(raw_text)
                        if cls._normalize_text(part)
                    ]
                for part in atomic_parts:
                    atomic_doc = Document(
                        page_content=part,
                        metadata=cls._copy_metadata(doc.metadata, raw_content=part),
                    )
                    windows.append([atomic_doc])
                continue

            if not current_window:
                current_window = [doc]
                current_chars = len(raw_text)
                current_sentences = sentence_count
                continue

            previous_doc = current_window[-1]
            previous_meta = previous_doc.metadata or {}
            current_meta = doc.metadata or {}
            same_page = previous_meta.get("page_number") == current_meta.get("page_number")
            next_chars = current_chars + 1 + len(raw_text)
            next_sentences = current_sentences + sentence_count

            should_flush = False
            if not same_page and current_chars >= cls.MIN_SEMANTIC_WINDOW_CHARS:
                should_flush = True
            elif next_chars > cls.MAX_SEMANTIC_WINDOW_CHARS:
                should_flush = True
            elif next_sentences > 6 and current_chars >= cls.MIN_SEMANTIC_WINDOW_CHARS:
                should_flush = True
            elif (
                len(current_window) >= cls.MAX_SEMANTIC_WINDOW_BLOCKS
                and current_chars >= cls.MIN_SEMANTIC_WINDOW_CHARS
            ):
                should_flush = True

            if should_flush:
                flush_window()
                current_window = [doc]
                current_chars = len(raw_text)
                current_sentences = sentence_count
                continue

            current_window.append(doc)
            current_chars = next_chars
            current_sentences = next_sentences

            if (
                current_chars >= cls.TARGET_SEMANTIC_WINDOW_CHARS
                and current_sentences >= 3
            ):
                flush_window()

        flush_window()
        return windows

    @classmethod
    def _make_semantic_window_doc(
        cls, docs: Sequence[Document], section_window_index: int
    ) -> Document | None:
        raw_parts = [cls._normalize_text(cls._source_content(doc)) for doc in docs]
        raw_parts = [part for part in raw_parts if part]
        if not raw_parts:
            return None

        metadata = dict((docs[0].metadata or {}))
        raw_content = cls._normalize_text(" ".join(raw_parts))
        if not cls._has_substantive_body(raw_content, metadata):
            return None

        metadata = cls._copy_metadata(
            metadata,
            raw_content=raw_content,
            chunk_strategy="semantic_window",
            chunk_stage="semantic_window",
            section_window_index=section_window_index,
            section_window_blocks=len(docs),
        )
        return Document(
            page_content=cls._semantic_chunk_text(raw_content, metadata),
            metadata=metadata,
        )

    @classmethod
    def _build_section_semantic_windows(cls, docs: Sequence[Document]) -> List[Document]:
        windows: List[Document] = []
        for idx, window_docs in enumerate(cls._semantic_seed_segments(docs), start=1):
            window_doc = cls._make_semantic_window_doc(window_docs, idx)
            if window_doc is not None:
                windows.append(window_doc)
        return windows

    @staticmethod
    def _cosine_similarity(vec_a: np.ndarray, vec_b: np.ndarray) -> float:
        denom = float(np.linalg.norm(vec_a) * np.linalg.norm(vec_b))
        if denom == 0:
            return 0.0
        return float(np.dot(vec_a, vec_b) / denom)

    def _semantic_merge_group(self, docs: List[Document]) -> List[Document]:
        if len(docs) <= 1:
            return docs

        model = self._get_semantic_model()
        embeddings = model.encode(
            [doc.page_content[:1200] for doc in docs],
            normalize_embeddings=True,
        )
        merged_docs: List[Document] = []
        current_embedding = np.array(embeddings[0])
        current_metadata = dict(docs[0].metadata)
        current_semantic_texts = [docs[0].page_content]
        current_raw_texts = [self._source_content(docs[0])]

        for idx in range(1, len(docs)):
            candidate_doc = docs[idx]
            candidate_embedding = np.array(embeddings[idx])
            candidate_raw_text = self._source_content(candidate_doc)
            similarity = self._cosine_similarity(current_embedding, candidate_embedding)
            merged_length = len(" ".join(current_raw_texts + [candidate_raw_text]))
            current_meta = current_metadata or {}
            candidate_meta = candidate_doc.metadata or {}
            structurally_adjacent = (
                current_meta.get("section_path") == candidate_meta.get("section_path")
                and current_meta.get("page_number") == candidate_meta.get("page_number")
                and candidate_meta.get("section_window_index")
                == current_meta.get("section_window_index", 0) + 1
            )
            small_window_merge = (
                len(self._normalize_text(" ".join(current_raw_texts))) < self.TARGET_SEMANTIC_WINDOW_CHARS
                or len(candidate_raw_text) < self.MIN_SEMANTIC_WINDOW_CHARS
            )

            if (
                structurally_adjacent
                and
                similarity >= self.SEMANTIC_SIMILARITY_THRESHOLD
                and small_window_merge
                and merged_length <= self.MAX_SEMANTIC_CHUNK_CHARS
            ):
                current_semantic_texts.append(candidate_doc.page_content)
                current_raw_texts.append(candidate_raw_text)
                current_embedding = np.mean(
                    np.stack([current_embedding, candidate_embedding]), axis=0
                )
                current_metadata = self._copy_metadata(
                    current_metadata, semantic_similarity=round(similarity, 4)
                )
                continue

            merged_raw = self._normalize_text(" ".join(current_raw_texts))
            merged_docs.append(
                Document(
                    page_content=self._semantic_chunk_text(merged_raw, current_metadata),
                    metadata=self._copy_metadata(
                        current_metadata,
                        raw_content=merged_raw,
                        chunk_strategy="semantic",
                        chunk_stage="semantic",
                        merged_from=len(current_semantic_texts),
                    ),
                )
            )
            current_embedding = candidate_embedding
            current_metadata = dict(candidate_doc.metadata)
            current_semantic_texts = [candidate_doc.page_content]
            current_raw_texts = [candidate_raw_text]

        merged_raw = self._normalize_text(" ".join(current_raw_texts))
        merged_docs.append(
            Document(
                page_content=self._semantic_chunk_text(merged_raw, current_metadata),
                metadata=self._copy_metadata(
                    current_metadata,
                    raw_content=merged_raw,
                    chunk_strategy="semantic",
                    chunk_stage="semantic",
                    merged_from=len(current_semantic_texts),
                    section_window_index=current_metadata.get("section_window_index", 1),
                ),
            )
        )
        return merged_docs

    @classmethod
    def _group_key(cls, doc: Document) -> Tuple[str, str, str]:
        metadata = doc.metadata or {}
        source = str(metadata.get("source", "unknown"))
        document_type = str(metadata.get("document_type", "unknown"))
        section_path = str(metadata.get("section_path", ""))
        page_number = metadata.get("page_number")
        page_key = f"page:{page_number}" if page_number is not None else "page:all"
        return (source, document_type, f"{page_key}|{section_path}")

    @classmethod
    def _prepare_retrievable_blocks(cls, documents: Sequence[Document]) -> List[Document]:
        prepared_docs: List[Document] = []
        grouped_docs: Dict[Tuple[str, str], List[Document]] = {}
        ordered_keys: List[Tuple[str, str]] = []

        for doc in documents:
            metadata = doc.metadata or {}
            key = (
                str(metadata.get("source", "unknown")),
                str(metadata.get("document_type", "unknown")),
            )
            if key not in grouped_docs:
                grouped_docs[key] = []
                ordered_keys.append(key)
            grouped_docs[key].append(doc)

        for key in ordered_keys:
            heading_stack: List[Tuple[int, str]] = []
            for doc in grouped_docs[key]:
                metadata = dict(doc.metadata or {})
                block_type = str(metadata.get("block_type", "paragraph"))
                if metadata.get("is_toc"):
                    continue
                if block_type == "heading":
                    level = int(metadata.get("heading_level", 1))
                    title = cls._normalize_text(doc.page_content)
                    heading_stack = [entry for entry in heading_stack if entry[0] < level]
                    heading_stack.append((level, title))
                    continue
                prepared_docs.append(
                    Document(
                        page_content=doc.page_content,
                        metadata=cls._attach_heading_context(metadata, heading_stack),
                    )
                )

        return prepared_docs

    @classmethod
    def _is_low_value_candidate(cls, doc: Document) -> bool:
        metadata = doc.metadata or {}
        if metadata.get("is_toc") or metadata.get("is_heading_only"):
            return True
        if metadata.get("block_type") == "table_row":
            return False
        return not cls._has_substantive_body(cls._source_content(doc), metadata)

    def split_docs(self, documents, chunk_size=1000, chunk_overlap=200):
        if not documents:
            return [], []

        prepared_docs = self._prepare_retrievable_blocks(documents)
        if not prepared_docs:
            return [], []

        grouped_prepared_docs: Dict[Tuple[str, str, str], List[Document]] = {}
        ordered_keys: List[Tuple[str, str, str]] = []
        for doc in prepared_docs:
            key = self._group_key(doc)
            if key not in grouped_prepared_docs:
                grouped_prepared_docs[key] = []
                ordered_keys.append(key)
            grouped_prepared_docs[key].append(doc)

        structured_chunks: List[Document] = []
        for key in ordered_keys:
            structured_chunks.extend(
                self._build_section_semantic_windows(grouped_prepared_docs[key])
            )
        if not structured_chunks:
            return [], []

        grouped_docs: Dict[Tuple[str, str, str], List[Document]] = {}
        ordered_chunk_keys: List[Tuple[str, str, str]] = []
        for doc in structured_chunks:
            key = self._group_key(doc)
            if key not in grouped_docs:
                grouped_docs[key] = []
                ordered_chunk_keys.append(key)
            grouped_docs[key].append(doc)

        parent_docs: List[Document] = []
        child_docs: List[Document] = []
        global_chunk_index = 0
        # Per-file counters so parent IDs are stable and collision-free per
        # source file — required for incremental add/remove of single files.
        per_file_counters: Dict[str, int] = {}
        for key in ordered_chunk_keys:
            for idx, doc in enumerate(self._semantic_merge_group(grouped_docs[key]), start=1):
                if self._is_low_value_candidate(doc):
                    continue
                global_chunk_index += 1
                file_key = self._file_key(str((doc.metadata or {}).get("source", "")))
                per_file_counters[file_key] = per_file_counters.get(file_key, 0) + 1
                metadata = self._copy_metadata(
                    doc.metadata,
                    chunk_id=f"{global_chunk_index}",
                    chunk_index=idx,
                    parent_doc_id=f"parent-{file_key}-{per_file_counters[file_key]}",
                    chunk_strategy="parent_section",
                    chunk_stage="parent_section",
                    parent_section=doc.metadata.get("section_path", ""),
                    content_preview=self._safe_preview(self._source_content(doc)),
                )
                if "merged_from" not in metadata:
                    metadata["merged_from"] = 1
                raw_content = self._source_content(doc)
                metadata.update(self._policy_topic_flags(raw_content, metadata))
                parent_doc = Document(
                    page_content=self._semantic_chunk_text(raw_content, metadata),
                    metadata=self._copy_metadata(metadata, raw_content=raw_content),
                )
                children = self._child_docs_from_parent(parent_doc)
                parent_doc = Document(
                    page_content=parent_doc.page_content,
                    metadata=self._copy_metadata(parent_doc.metadata, child_count=len(children)),
                )
                parent_docs.append(parent_doc)
                child_docs.extend(children)

        return child_docs, parent_docs
